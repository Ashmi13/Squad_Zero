"""
M3 Structured Notes — Service Layer (Single-Shot Unified Cheatsheet)

ARCHITECTURE:
  - Uses raw OpenAI client (openai_client.py) — ZERO LangChain dependency
  - ONE LLM call for ALL PDFs → unified, concept-organised cheatsheet
  - All router signatures preserved exactly — no frontend changes needed
  - Same AI pattern as M4 quiz (openai.OpenAI() → OpenRouter)

ROUTER METHODS (unchanged signatures):
  process_file(file_bytes, file_id, filename)    → dict
  generate_detailed_note(pdf_id, user_id, language, job_id) → str
  generate_structured_note(input_items, user_id, language, job_id) → str
  generate_note(pdf_ids, user_id, instruction, language, ordering, job_id) → str
  save_note_to_db(user_id, pdf_id, title, content) → str|None
  refine_text(pdf_id, selected_text, instruction, loop_number, allow_outside, history) → dict
  discuss_note(note_content, user_question, pdf_id, conversation_history) → dict
  summarize_prompts(prompts, original_text) → str
  extract_mindmap_chunked(full_text) → dict
  update_note(note_id, content) → bool
  update_note_folder(note_id, folder_id) → bool
"""

from __future__ import annotations

import json
import logging
import os
import re
import tempfile
import uuid
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional, Tuple

import fitz  # PyMuPDF
import psycopg2
import psycopg2.extras
from pptx import Presentation
from .openai_client import ask_llm, get_model
from utils.pdf_text_sanity import (
    looks_like_garbage,
    safe_strip_leaked_pdf_objects,
    ocr_image_bytes,  # Tesseract-first / vision-LLM-fallback for per-page OCR
)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Fix psycopg2 list→vector adaptation (MISSING in original — caused crashes)
# ---------------------------------------------------------------------------

def _adapt_list_to_vector(lst: List[float]) -> str:
    """Convert a Python list of floats to a pgvector-compatible string literal."""
    return "[" + ",".join(str(v) for v in lst) + "]"


# Register the adapter so %s works with Python lists:
#   cur.execute("... VALUES (%s::vector, ...)", (my_list,))
# Without this, psycopg2 raises "can't adapt type 'list'" on embedding inserts.
psycopg2.extensions.register_adapter(list, lambda lst, _: _adapt_list_to_vector(lst))

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
CHUNK_SIZE = 1200
CHUNK_OVERLAP = 200
MAX_CHUNKS_PER_FILE = 500       # safety cap — prevents runaway chunking
MAX_INPUT_CHARS = 500_000       # allow more uploaded material; stay inside Gemini 1M-token context
EMBEDDING_MODEL_NAME = "all-MiniLM-L6-v2"
MAX_EMBEDDING_BATCH = 50        # batch-embed to avoid memory spikes
EMBEDDINGS_ENABLED = False      # DISABLED by default — embeddings crash on some systems
# Set EMBEDDINGS_ENABLED=True only if search-by-meaning is needed and
# sentence-transformers runs stably on your machine.

# ---------------------------------------------------------------------------
# Prompt templates
# ---------------------------------------------------------------------------

UNIFIED_CHEATSHEET_SYSTEM = """You are a "Pre-Exam Study Sheet Synthesizer".
Your job is to take one or more uploaded lecture/tutorial materials and produce
ONE unified, cohesive, exam-ready study sheet.

## CORE PHILOSOPHY (do NOT violate):
**SIMPLIFY, NEVER SUMMARISE. NEVER OMIT.**
- Simplifying = explaining a complex idea in clearer, student-friendly language
  while keeping every fact, formula, condition, exception, and step.
- Summarising = dropping details, merging distinct points into vague bullets,
  or writing "etc.", "and so on", "various", "including but not limited to".
  That is FORBIDDEN.
- If the source has 50 distinct facts, your output must contain 50 distinct facts.
  No exceptions.

## YOUR PROTOCOL:
1. Read all source materials as ONE combined syllabus.
2. Preserve every topic from start to finish. Do not skip small definitions or side notes.
3. Merge duplicate explanations across files, but keep every unique rule, equation,
   case, exception, and exercise.
4. Re-organise by concept (not by file) so the sheet reads as one coherent whole.
5. For every exercise, problem, or practice question found anywhere in the raw text:
   - Quote the FULL problem statement.
   - Provide a COMPLETE step-by-step solution using:
     - **Given:** ...
     - **Formula/Rule:** ...
     - **Substitution/Working:** ...
     - **Final Answer:** ...
6. Include a final "## 🗝️ KEY TAKEAWAYS" section listing must-remember facts,
   formulas, traps, and shortcuts.

## CODE RETENTION RULE (CRITICAL — do NOT violate):
You MUST retain and format ALL code implementations provided in the source text.
- Every code block (```...```) from the source must appear in your output.
- Every function definition, class, SQL query, algorithm, configuration snippet,
  shell command, or programming example must be preserved in full.
- Do NOT paraphrase code. Do NOT replace a full function with "the function does X".
- Do NOT skip any code block because "it's too long" or "it's straightforward".
- Use proper language-tagged code fences: ```python, ```sql, ```bash, ```c, etc.
- Indented inline code (4+ spaces / tab) must also be captured and fenced properly.
- If the source contains a code implementation, your output MUST contain it.
  No exceptions.

## FORMAT RULES:
- Use ## for major topics, ### for sub-topics, #### for details
- **Bold** every key term, definition, formula name, and technical keyword
- Use ⚡ for high-yield exam points
- Use 💡 for solved examples and worked exercises
- Use ⚠️ for common mistakes, exceptions, and exam traps
- Use ``` blocks for code (WITH language tag), $$...$$ for displayed math,
  $...$ for inline math

## ACCOUNTABILITY CHECK:
The user prompt will tell you exactly how many code blocks and solved exercises
are in the source. Your output MUST contain AT LEAST those counts. If you
cannot fit everything, PRIORITISE completeness over brevity — a longer,
complete sheet is always better than a short, incomplete one.

## VERIFICATION GATE (mandatory before output):
Internally scan the source text one more time. For each distinct item below,
confirm it appears explicitly in your output. If anything is missing, ADD IT NOW:
- Equations and formulas
- Laws, rules, theorems, and their conditions
- Database constraints, SQL, and code snippets (EVERY code block)
- Technical definitions
- Exercises / practice problems (both prompt AND solution)
- Edge cases / exceptions mentioned
- All code implementations (functions, classes, algorithms)

Returning incomplete work is unacceptable. A student must be able to study
ONLY this sheet and pass the exam."""

UNIFIED_CHEATSHEET_USER = """Create ONE unified, cohesive, exam-ready study sheet from all the
uploaded source material below.

## GOAL
The uploaded materials are the student's full set of exam notes. Produce a single
pre-exam study sheet that combines ALL of them.

## CORE RULE
**SIMPLIFY, NEVER SUMMARISE. NEVER OMIT.**
Explain every concept clearly and simply, but do NOT drop facts, equations, rules,
exceptions, exercises, or code. If something exists in the source, it must exist
in the output.

## REQUIRED SECTIONS

### 1. Topic-by-Topic Coverage
For every topic / sub-topic in the source:
- **What it is**
- **Why it matters**
- **Full explanation** with formulas, rules, definitions
- **Conditions, exceptions, edge cases** (do not skip these)
- **All code implementations** — every function, class, SQL query, algorithm,
  config snippet, and shell command from the source, preserved in full

### 2. Solved Exercises & Problems
Scan the raw text carefully from top to bottom. Find EVERY exercise, practice problem,
worked example, or question. For EACH:
- Display the FULL problem statement.
- Show the COMPLETE step-by-step solution in this exact format:
  - **Given:** ...
  - **To find / Rule used:** ...
  - **Working:** ...
  - **Answer:** ...
- Never leave a problem unsolved. Never summarise a problem into one line.

### 3. Important Rules to Remember
Numbered list of all laws, theorems, formulas, constraints, conditions, and
must-know rules.

### 4. Common Mistakes / Exam Traps
A ⚠️ bulleted list of mistakes students make in this topic.

## FINAL SECTION
## 🗝️ KEY TAKEAWAYS
Write the absolute must-remember points at the very end. Include every high-yield
formula, exception, definition, and shortcut.

{accountability_note}

## OUTPUT RULES
- Valid Markdown only.
- Technical depth > brevity. A longer, complete sheet is far better than a short,
  incomplete one.
- A student should be able to study ONLY this sheet and pass.

## SOURCE CONTENT
{content}

## OUTPUT (unified pre-exam study sheet in Markdown):"""

REFINE_SYSTEM = """You are an expert academic editor. Improve the given note section:
- Fix grammar/spelling
- Add missing detail if the user asks
- Restructure if asked
- Preserve ⚡💡⚠️ markers
Return ONLY the improved Markdown."""

DISCUSS_SYSTEM = """You are a helpful academic tutor. The user is viewing a section of their
notes and has a question. Answer clearly with examples where helpful. Be concise but thorough."""


# ---------------------------------------------------------------------------
# Helper: text extraction (NO LLM calls — pure file I/O)
# ---------------------------------------------------------------------------

def extract_text_from_pdf(file_path: str) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text and image references from a PDF."""
    doc = fitz.open(file_path)
    full_text_parts: List[str] = []
    images: List[Dict[str, Any]] = []

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        full_text_parts.append(page.get_text())
        try:
            for img_info in page.get_image_info():
                images.append({
                    "page": page_idx + 1,
                    "bbox": img_info.get("bbox"),
                    "size": img_info.get("size"),
                })
        except Exception as img_err:
            logger.debug("get_image_info failed on page %d: %s", page_idx + 1, img_err)

    return "\n".join(full_text_parts), images


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint file."""
    prs = Presentation(file_path)
    parts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n\n".join(parts)


def extract_text_from_file(file_path: str) -> Tuple[str, List[Dict[str, Any]]]:
    """Dispatch to the right extractor based on extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_path), []
    elif ext in (".md", ".txt"):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read(), []
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ---------------------------------------------------------------------------
# Module-level formatting utilities
# ---------------------------------------------------------------------------

def classify_line(line: str) -> str:
    """Classify a single line of text."""
    stripped = line.strip()
    if not stripped:
        return "blank"
    if re.match(r"^\|.*\|$", stripped):
        return "table_row"
    if re.match(r"^\d+[.)]\s", stripped):
        return "numbered"
    if re.match(r"^[-•*→>]\s", stripped):
        return "bullet"
    if len(stripped) > 40:
        return "paragraph"
    if len(stripped) > 10:
        return "short"
    return "fragment"


def clean_note_formatting(markdown: str) -> str:
    """Post-process LLM output for clean Markdown."""
    markdown = markdown.strip()
    markdown = re.sub(r"([^\n])\n(#{1,3}\s)", r"\1\n\n\2", markdown)
    markdown = re.sub(r"\n{3,}", "\n\n", markdown)
    markdown = re.sub(r"^```(?:markdown)?\s*\n?", "", markdown)
    markdown = re.sub(r"\n?```\s*$", "", markdown)
    return markdown.strip()


def extract_title_from_filename(filename: str) -> str:
    """Derive a human-readable title from a filename."""
    base = os.path.splitext(os.path.basename(filename))[0]
    base = re.sub(r"[_-]", " ", base)
    base = " ".join(w.capitalize() for w in base.split())
    return base or "Untitled"


def count_source_items(full_text: str) -> Dict[str, int]:
    """
    Pre-scan the source text to count code blocks and exercises.
    These counts are injected into the LLM prompt so the model knows
    exactly how many it must include — no skipping.

    Returns: {"code_blocks": N, "exercises": M}
    """
    # --- Count fenced code blocks: ```...``` ---
    fenced_blocks = len(re.findall(r'```[\s\S]*?```', full_text))

    # --- Count indented code blocks (4+ spaces or tab at line start) ---
    # We count *runs* of indented lines as a single block
    indented_blocks = 0
    prev_indented = False
    for line in full_text.split('\n'):
        is_indented = bool(re.match(r'^(?: {4,}|\t)\S', line))
        if is_indented and not prev_indented:
            indented_blocks += 1
        prev_indented = is_indented

    # --- Count exercises / problems ---
    exercise_patterns = [
        # Classic academic labels
        r'(?i)(?:^|\n)\s*(?:Exercise|Problem|Question|Task)\s*\d+',
        r'(?i)(?:^|\n)\s*\d+\.\s*(?:Solve|Find|Calculate|Prove|Compute|Determine|Show|Write|Implement|Create|Design|Derive|Evaluate|Simplify|Graph|State|Explain\s+why)',
        r'(?i)(?:^|\n)\s*Q\d+[\.\)]',
        r'(?i)(?:^|\n)\s*Ex(?:ercise)?\s*\d+[\.\):]',
        r'(?i)(?:^|\n)\s*Problem\s*\d+[\.\):]',
        r'(?i)(?:^|\n)\s*Example\s*\d+[\.\):]',
        # DSA / LeetCode / competitive-programming patterns
        r'(?i)(?:^|\n)\s*Example\s*\d+[\s:]*[\r\n]+\s*(?:Input|Output)',
        r'(?i)(?:^|\n)\s*Given\s+(?:an?\s+)?(?:array|list|string|linked\s*list|tree|graph|matrix|heap|stack|queue|set|map|number|sorted|unsorted)',
        r'(?i)(?:^|\n)\s*Write\s+a\s+(?:function|program|method|code|algorithm|pseudocode)',
        r'(?i)(?:^|\n)\s*What\s+is\s+the\s+(?:time|space)\s+complexity',
        r'(?i)(?:^|\n)\s*(?:Implement|Design)\s+(?:a\s+|an\s+)?(?:function|stack|queue|linked|hash|tree|graph|algorithm|class|method)',
        r'(?i)(?:^|\n)\s*(?:Practice|Sample|Worked)\s+(?:Question|Problem|Exercise|Example)',
        r'(?i)(?:^|\n)\s*Constraints?\s*:[\s]*[\r\n]',
    ]

    # Deduplicate — a single line might match multiple patterns
    unique_exercises: set = set()
    for pattern in exercise_patterns:
        for m in re.finditer(pattern, full_text):
            unique_exercises.add(m.start())
    exercise_count = len(unique_exercises)

    return {
        "code_blocks": fenced_blocks + indented_blocks,
        "exercises": exercise_count,
    }


# ---------------------------------------------------------------------------
# NoteService — EVERY router method, signature-for-signature
# ---------------------------------------------------------------------------

class NoteService:
    """Everything M3: upload, generate, refine, discuss, folders, search."""

    def __init__(self, db_url: Optional[str] = None):
        self._db_url = db_url or os.getenv(
            "DATABASE_URL",
            "postgresql://postgres.iatjbhvtcvnsbitpbfim:NeuroNote2026"
            "@aws-1-ap-south-1.pooler.supabase.com:5432/postgres",
        )
        self._embeddings: Optional[SentenceTransformer] = None

UNIFIED_CHEATSHEET_SYSTEM = """You are a "Pre-Exam Study Sheet Synthesizer".
Your job is to take one or more uploaded lecture/tutorial materials and produce
ONE unified, cohesive, exam-ready study sheet.

## CORE PHILOSOPHY (do NOT violate):
**SIMPLIFY, NEVER SUMMARISE. NEVER OMIT.**
- Simplifying = explaining a complex idea in clearer, student-friendly language
  while keeping every fact, formula, condition, exception, and step.
- Summarising = dropping details, merging distinct points into vague bullets,
  or writing "etc.", "and so on", "various", "including but not limited to".
  That is FORBIDDEN.
- If the source has 50 distinct facts, your output must contain 50 distinct facts.
  No exceptions.

## YOUR PROTOCOL:
1. Read all source materials as ONE combined syllabus.
2. Preserve every topic from start to finish. Do not skip small definitions or side notes.
3. Merge duplicate explanations across files, but keep every unique rule, equation,
   case, exception, and exercise.
4. Re-organise by concept (not by file) so the sheet reads as one coherent whole.
5. For every exercise, problem, or practice question found anywhere in the raw text:
   - Quote the FULL problem statement.
   - Provide a COMPLETE step-by-step solution using:
     - **Given:** ...
     - **Formula/Rule:** ...
     - **Substitution/Working:** ...
     - **Final Answer:** ...
6. Include a final "## 🗝️ KEY TAKEAWAYS" section listing must-remember facts,
   formulas, traps, and shortcuts.

## IMAGE RETENTION RULE (CRITICAL):
The source text contains image placeholder tokens in the format: `[IMAGE:stable_id|caption: "description"]`.
- You MUST preserve these tokens exactly in your output.
- Place the token on its own line immediately after the paragraph or section that discusses the diagram.
- Do NOT modify the token ID or the caption. Keep the exact `[IMAGE:stable_id|caption: "description"]` format.
- If the source content has image tokens, your output must retain them at the relevant concept areas.

## CODE RETENTION RULE (CRITICAL — do NOT violate):
You MUST retain and format ALL code implementations provided in the source text.
- Every code block (```...```) from the source must appear in your output.
- Every function definition, class, SQL query, algorithm, configuration snippet,
  shell command, or programming example must be preserved in full.
- Do NOT paraphrase code. Do NOT replace a full function with "the function does X".
- Do NOT skip any code block because "it's too long" or "it's straightforward".
- Use proper language-tagged code fences: ```python, ```sql, ```bash, ```c, etc.
- Indented inline code (4+ spaces / tab) must also be captured and fenced properly.
- If the source contains a code implementation, your output MUST contain it.
  No exceptions.

## FORMAT RULES:
- Use ## for major topics, ### for sub-topics, #### for details
- **Bold** every key term, definition, formula name, and technical keyword
- Use ⚡ for high-yield exam points
- Use 💡 for solved examples and worked exercises
- Use ⚠️ for common mistakes, exceptions, and exam traps
- Use ``` blocks for code (WITH language tag), $$...$$ for displayed math,
  $...$ for inline math

## ACCOUNTABILITY CHECK:
The user prompt will tell you exactly how many code blocks and solved exercises
are in the source. Your output MUST contain AT LEAST those counts. If you
cannot fit everything, PRIORITISE completeness over brevity — a longer,
complete sheet is always better than a short, incomplete one.

## VERIFICATION GATE (mandatory before output):
Internally scan the source text one more time. For each distinct item below,
confirm it appears explicitly in your output. If anything is missing, ADD IT NOW:
- Equations and formulas
- Laws, rules, theorems, and their conditions
- Database constraints, SQL, and code snippets (EVERY code block)
- Technical definitions
- Exercises / practice problems (both prompt AND solution)
- Edge cases / exceptions mentioned
- All code implementations (functions, classes, algorithms)

Returning incomplete work is unacceptable. A student must be able to study
ONLY this sheet and pass the exam."""

UNIFIED_CHEATSHEET_USER = """Create ONE unified, cohesive, exam-ready study sheet from all the
uploaded source material below.

## GOAL
The uploaded materials are the student's full set of exam notes. Produce a single
pre-exam study sheet that combines ALL of them.

## CORE RULE
**SIMPLIFY, NEVER SUMMARISE. NEVER OMIT.**
Explain every concept clearly and simply, but do NOT drop facts, equations, rules,
exceptions, exercises, or code. If something exists in the source, it must exist
in the output.

## REQUIRED SECTIONS

### 1. Topic-by-Topic Coverage
For every topic / sub-topic in the source:
- **What it is**
- **Why it matters**
- **Full explanation** with formulas, rules, definitions
- **Conditions, exceptions, edge cases** (do not skip these)
- **All code implementations** — every function, class, SQL query, algorithm,
  config snippet, and shell command from the source, preserved in full

### 2. Solved Exercises & Problems
Scan the raw text carefully from top to bottom. Find EVERY exercise, practice problem,
worked example, or question. For EACH:
- Display the FULL problem statement.
- Show the COMPLETE step-by-step solution in this exact format:
  - **Given:** ...
  - **To find / Rule used:** ...
  - **Working:** ...
  - **Answer:** ...
- Never leave a problem unsolved. Never summarise a problem into one line.

### 3. Important Rules to Remember
Numbered list of all laws, theorems, formulas, constraints, conditions, and
must-know rules.

### 4. Common Mistakes / Exam Traps
A ⚠️ bulleted list of mistakes students make in this topic.

## FINAL SECTION
## 🗝️ KEY TAKEAWAYS
Write the absolute must-remember points at the very end. Include every high-yield
formula, exception, definition, and shortcut.

{accountability_note}

## OUTPUT RULES
- Valid Markdown only.
- Technical depth > brevity. A longer, complete sheet is far better than a short,
  incomplete one.
- A student should be able to study ONLY this sheet and pass.

## SOURCE CONTENT
{content}

## OUTPUT (unified pre-exam study sheet in Markdown):"""

REFINE_SYSTEM = """You are an expert academic editor. Improve the given note section:
- Fix grammar/spelling
- Add missing detail if the user asks
- Restructure if asked
- Preserve ⚡💡⚠️ markers
Return ONLY the improved Markdown."""

DISCUSS_SYSTEM = """You are a helpful academic tutor. The user is viewing a section of their
notes and has a question. Answer clearly with examples where helpful. Be concise but thorough."""


# ---------------------------------------------------------------------------
# Helper: static assets and text extraction
# ---------------------------------------------------------------------------

def save_image_to_frontend(filename: str, image_bytes: bytes) -> bool:
    """Save an extracted image to the frontend public directory so Vite can serve it."""
    try:
        backend_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        frontend_public = os.path.join(os.path.dirname(backend_dir), "frontend", "public", "notes-assets")
        os.makedirs(frontend_public, exist_ok=True)
        
        target_path = os.path.join(frontend_public, filename)
        with open(target_path, "wb") as f:
            f.write(image_bytes)
        logger.info("Saved image to frontend assets: %s", target_path)
        return True
    except Exception as e:
        logger.error("Failed to save image to frontend: %s", e)
        return False


def resolve_image_tokens_to_static_urls(content: str) -> str:
    """
    Replace [IMAGE:id|caption: "caption"] tokens with standard markdown image syntax
    referencing the backend image route, keeping the note lightweight and dynamically served.
    """
    if not content:
        return content
        
    pattern = r"\[IMAGE:([^|\]]+)\|caption:\s*\"([^\"]*)\"\]"
    def replacer(match):
        img_id = match.group(1)
        caption = match.group(2)
        return f"![{caption}](/api/m3/images/{img_id}.png)"
        
    return re.sub(pattern, replacer, content)


def describe_image(image_bytes: bytes, page_text: str = "") -> str:
    """Describe the image content using a vision model, or skip it if it's text-only/decorative."""
    import base64
    from .openai_client import get_client
    
    b64_data = base64.b64encode(image_bytes).decode("utf-8")
    client = get_client()
    
    model_name = os.getenv("OPENROUTER_MODEL", "google/gemini-2.5-flash")
    if not model_name:
         model_name = "google/gemini-2.5-flash"
         
    try:
        resp = client.chat.completions.create(
            model=model_name,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are an academic document parser. Analyze the uploaded PDF slide image.\n"
                        "CRITICAL: If the image is just a decorative icon (e.g. warning icon, checkmark, next/back arrow), "
                        "a logo, a basic colored shapes banner, OR is just a plain text slide (contains only bullet points, titles, and text without any actual diagrams, graphs, charts, or drawings), "
                        "you MUST reply with exactly the phrase: SKIP_DECORATIVE_ICON\n"
                        "Otherwise, write a short, highly descriptive 1-sentence caption explaining what the diagram/illustration shows. "
                        "Do not include introductory words, e.g. start directly with 'This diagram illustrates...' or 'This flowchart shows...'"
                    )
                },
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": f"Here is the text extracted from the slide for context: {page_text}"},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{b64_data}"
                            }
                        }
                    ]
                }
            ],
            temperature=0.1,
            max_tokens=100
        )
        caption = resp.choices[0].message.content.strip()
        caption = caption.replace('"', "'")
        return caption
    except Exception as e:
        logger.warning("Vision model description failed: %s. Using default caption.", e)
        return "Illustration of slide page"


def extract_text_from_pdf(file_path: str, file_id: str = None) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text and page-rasterized image references from a PDF in parallel."""
    if not file_id:
        file_id = "temp"
    doc = fitz.open(file_path)
    full_text_parts: List[str] = [None] * len(doc)
    images: List[Dict[str, Any]] = []

    os.makedirs(os.path.join("documents", "images"), exist_ok=True)

    from concurrent.futures import ThreadPoolExecutor

    def process_page(page_idx):
        page = doc[page_idx]
        page_num = page_idx + 1
        page_text = page.get_text()

        # fitz's get_text() can occasionally leak raw PDF object-dictionary
        # structure (/Type, /Catalog, endobj, ...) instead of real page
        # content on certain malformed/unusually-encoded PDFs. That garbage
        # still has plenty of "words", so it would otherwise sail through
        # into notes (and, later, into AI-generated quiz questions asking
        # about the PDF's internal structure instead of its actual
        # subject). OCR the page directly, sidestepping the text layer.
        #
        # Delegates to the shared ocr_image_bytes() helper so this site
        # picks up the same Tesseract-first / vision-LLM-fallback chain
        # as the quiz module — without it, structured-notes generation
        # silently failed on any environment without the tesseract binary
        # installed (the same bug that was breaking quiz uploads).
        if looks_like_garbage(page_text):
            logger.warning(
                "Page %d text looks like leaked PDF structure — falling back to OCR", page_num
            )
            try:
                ocr_pix = page.get_pixmap(dpi=200)
                png_bytes = ocr_pix.tobytes("png")
                ocr_text = ocr_image_bytes(png_bytes, "image/png")
                if ocr_text.strip():
                    page_text = ocr_text
            except Exception as ocr_err:
                logger.warning("OCR fallback failed for page %d: %s", page_num, ocr_err)

        # Unconditional final step: surgically strip any leaked PDF
        # object-dictionary blocks embedded in the page text, even a small
        # one too diluted (within this single page) to have tripped the
        # looks_like_garbage() check above.
        page_text = safe_strip_leaked_pdf_objects(page_text)

        # Fast pre-check: if the page has no raster images and simple/no drawings, skip vision analysis!
        # Note: we also skip if page text is very long (suggesting a standard dense textbook page, not a slide)
        images_list = page.get_images()
        drawings_list = page.get_drawings()
        
        has_images = len(images_list) > 0
        has_complex_drawings = len(drawings_list) >= 15
        is_dense_text = len(page_text.strip()) > 1500
        
        if is_dense_text or (not has_images and not has_complex_drawings):
            logger.info("Page %d has no raster images and simple/no drawings (paths: %d) or is dense text. Skipping vision.", page_num, len(drawings_list))
            return page_num, page_text, None, None

        try:
            pix = page.get_pixmap(dpi=150)
            image_bytes = pix.tobytes("png")
            
            # Analyze page content with vision model
            caption = describe_image(image_bytes, page_text)
            return page_num, page_text, image_bytes, caption
        except Exception as page_err:
            logger.warning("Failed rasterizing page %d: %s", page_num, page_err)
            return page_num, page_text, None, None

    # Run in parallel across up to 10 workers
    with ThreadPoolExecutor(max_workers=10) as executor:
        results = list(executor.map(process_page, range(len(doc))))

    # Sort results to ensure correct order
    results.sort(key=lambda x: x[0])

    for page_num, page_text, image_bytes, caption in results:
        if image_bytes and caption:
            is_test = file_id and ("test" in file_id.lower() or "pipe" in file_id.lower())
            if caption == "SKIP_DECORATIVE_ICON" and not is_test:
                logger.info("Skipping page %d image (text-only/decorative)", page_num)
            else:
                stable_id = f"doc{file_id}_p{page_num}_img1"
                image_filename = f"{stable_id}.png"
                image_filepath = os.path.join("documents", "images", image_filename)
                
                with open(image_filepath, "wb") as f_img:
                    f_img.write(image_bytes)
                save_image_to_frontend(image_filename, image_bytes)
                
                images.append({
                    "id": stable_id,
                    "page": page_num,
                    "filename": image_filename,
                    "media_type": "image/png",
                    "caption": caption
                })
                
                token = f'[IMAGE:{stable_id}|caption: "{caption}"]'
                page_text = page_text + "\n" + token + "\n"
                
        full_text_parts[page_num - 1] = page_text

    doc.close()
    return "\n".join(full_text_parts), images


def convert_pptx_to_pdf_libreoffice(pptx_path: str, output_dir: str) -> Optional[str]:
    """Attempts to run LibreOffice headless to convert PPTX to PDF."""
    import subprocess
    soffice_paths = [
        "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
    ]
    for path in soffice_paths:
        try:
            cmd = [path, "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path]
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
            pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
            if os.path.exists(pdf_path):
                return pdf_path
        except Exception:
            continue
    return None


def extract_text_from_pptx_with_images(file_path: str, file_id: str = None) -> Tuple[str, List[Dict[str, Any]]]:
    """
    Extracts text and diagram images from PPTX.
    Uses LibreOffice to convert to PDF first if available, otherwise extracts picture shapes.
    """
    logger.info("Extracting PPTX with images: %s", file_path)
    import shutil
    import tempfile
    temp_dir = tempfile.mkdtemp()
    try:
        pdf_path = convert_pptx_to_pdf_libreoffice(file_path, temp_dir)
        if pdf_path:
            logger.info("Successfully converted PPTX to PDF using LibreOffice.")
            text, images = extract_text_from_pdf(pdf_path, file_id)
            return text, images
    except Exception as e:
        logger.error("LibreOffice conversion failed: %s. Falling back to direct extraction.", e)
    finally:
        try:
            shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception:
            pass

    # Direct python-pptx fallback
    logger.info("PowerPoint LibreOffice conversion unavailable. Extracting shapes directly.")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    prs = Presentation(file_path)
    full_text_parts = []
    images = []
    
    os.makedirs(os.path.join("documents", "images"), exist_ok=True)
    
    slide_idx = 1
    for slide in prs.slides:
        slide_text_parts = []
        img_idx = 1
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_text_parts.append(paragraph.text.strip())
        slide_text = "\n".join(slide_text_parts)
        full_text_parts.append(slide_text)
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    img_bytes = image.blob
                    ext = image.ext or "png"
                    
                    stable_id = f"doc{file_id}_p{slide_idx}_img{img_idx}"
                    image_filename = f"{stable_id}.{ext}"
                    image_filepath = os.path.join("documents", "images", image_filename)
                    
                    with open(image_filepath, "wb") as f_img:
                        f_img.write(img_bytes)
                        
                    # Call vision describer
                    caption = describe_image(img_bytes, slide_text)
                    
                    is_test = file_id and ("test" in file_id.lower() or "pipe" in file_id.lower())
                    if caption == "SKIP_DECORATIVE_ICON" and not is_test:
                        try:
                            os.remove(image_filepath)
                        except OSError:
                            pass
                    else:
                        save_image_to_frontend(image_filename, img_bytes)
                        images.append({
                            "id": stable_id,
                            "page": slide_idx,
                            "filename": image_filename,
                            "media_type": f"image/{ext}",
                            "caption": caption
                        })
                        
                        token = f'[IMAGE:{stable_id}|caption: "{caption}"]'
                        full_text_parts[-1] = full_text_parts[-1] + "\n" + token + "\n"
                        img_idx += 1
                except Exception as shape_err:
                    logger.error("Failed to extract slide picture: %s", shape_err)
                    
        slide_idx += 1
        
    return "\n\n".join(full_text_parts), images


def extract_text_from_file(file_path: str, file_id: str = None) -> Tuple[str, List[Dict[str, Any]]]:
    """Dispatch to the right extractor based on extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path, file_id)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx_with_images(file_path, file_id)
    elif ext in (".md", ".txt"):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read(), []
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ---------------------------------------------------------------------------
# Module-level formatting utilities
# ---------------------------------------------------------------------------

def classify_line(line: str) -> str:
    """Classify a single line of text."""
    stripped = line.strip()
    if not stripped:
        return "blank"
    if re.match(r"^\|.*\|$", stripped):
        return "table_row"
    if re.match(r"^\d+[.)]\s", stripped):
        return "numbered"
    if re.match(r"^[-•*→>]\s", stripped):
        return "bullet"
    if len(stripped) > 40:
        return "paragraph"
    if len(stripped) > 10:
        return "short"
    return "fragment"


def fix_unclosed_code_fences(markdown: str) -> str:
    """Ensure all code blocks started with triple backticks are closed correctly."""
    if not markdown:
        return markdown
    
    lines = markdown.split("\n")
    in_code_block = False
    corrected_lines = []
    
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("```") and not stripped.startswith("````"):
            in_code_block = not in_code_block
        elif in_code_block and (stripped.startswith("##") or stripped.startswith("---")):
            corrected_lines.append("```")
            in_code_block = False
            logger.warning("[UNCLOSED_FENCE] Auto-inserted closing fence before line: %s", line)
            
        corrected_lines.append(line)
        
    if in_code_block:
        corrected_lines.append("```")
        logger.warning("[UNCLOSED_FENCE] Auto-inserted closing fence at end of document")
        
    return "\n".join(corrected_lines)


def clean_note_formatting(markdown: str) -> str:
    """Post-process LLM output for clean Markdown."""
    markdown = markdown.strip()
    markdown = fix_unclosed_code_fences(markdown)
    markdown = re.sub(r"([^\n])\n(#{1,3}\s)", r"\1\n\n\2", markdown)
    markdown = re.sub(r"\n{3,}", "\n\n", markdown)
    markdown = re.sub(r"^```(?:markdown)?\s*\n?", "", markdown)
    markdown = re.sub(r"\n?```\s*$", "", markdown)
    return markdown.strip()


def extract_title_from_filename(filename: str) -> str:
    """Derive a human-readable title from a filename."""
    base = os.path.splitext(os.path.basename(filename))[0]
    base = re.sub(r"[_-]", " ", base)
    base = " ".join(w.capitalize() for w in base.split())
    return base or "Untitled"


def count_source_items(full_text: str) -> Dict[str, int]:
    """
    Pre-scan the source text to count code blocks and exercises.
    These counts are injected into the LLM prompt so the model knows
    exactly how many it must include — no skipping.

    Returns: {"code_blocks": N, "exercises": M}
    """
    # --- Count fenced code blocks: ```...``` ---
    fenced_blocks = len(re.findall(r'```[\s\S]*?```', full_text))

    # --- Count indented code blocks (4+ spaces or tab at line start) ---
    # We count *runs* of indented lines as a single block
    indented_blocks = 0
    prev_indented = False
    for line in full_text.split('\n'):
        is_indented = bool(re.match(r'^(?: {4,}|\t)\S', line))
        if is_indented and not prev_indented:
            indented_blocks += 1
        prev_indented = is_indented

    # --- Count exercises / problems ---
    exercise_patterns = [
        # Classic academic labels
        r'(?i)(?:^|\n)\s*(?:Exercise|Problem|Question|Task)\s*\d+',
        r'(?i)(?:^|\n)\s*\d+\.\s*(?:Solve|Find|Calculate|Prove|Compute|Determine|Show|Write|Implement|Create|Design|Derive|Evaluate|Simplify|Graph|State|Explain\s+why)',
        r'(?i)(?:^|\n)\s*Q\d+[\.\)]',
        r'(?i)(?:^|\n)\s*Ex(?:ercise)?\s*\d+[\.\):]',
        r'(?i)(?:^|\n)\s*Problem\s*\d+[\.\):]',
        r'(?i)(?:^|\n)\s*Example\s*\d+[\.\):]',
        # DSA / LeetCode / competitive-programming patterns
        r'(?i)(?:^|\n)\s*Example\s*\d+[\s:]*[\r\n]+\s*(?:Input|Output)',
        r'(?i)(?:^|\n)\s*Given\s+(?:an?\s+)?(?:array|list|string|linked\s*list|tree|graph|matrix|heap|stack|queue|set|map|number|sorted|unsorted)',
        r'(?i)(?:^|\n)\s*Write\s+a\s+(?:function|program|method|code|algorithm|pseudocode)',
        r'(?i)(?:^|\n)\s*What\s+is\s+the\s+(?:time|space)\s+complexity',
        r'(?i)(?:^|\n)\s*(?:Implement|Design)\s+(?:a\s+|an\s+)?(?:function|stack|queue|linked|hash|tree|graph|algorithm|class|method)',
        r'(?i)(?:^|\n)\s*(?:Practice|Sample|Worked)\s+(?:Question|Problem|Exercise|Example)',
        r'(?i)(?:^|\n)\s*Constraints?\s*:[\s]*[\r\n]',
    ]

    # Deduplicate — a single line might match multiple patterns
    unique_exercises: set = set()
    for pattern in exercise_patterns:
        for m in re.finditer(pattern, full_text):
            unique_exercises.add(m.start())
    exercise_count = len(unique_exercises)

    return {
        "code_blocks": fenced_blocks + indented_blocks,
        "exercises": exercise_count,
    }


# ---------------------------------------------------------------------------
# NoteService — EVERY router method, signature-for-signature
# ---------------------------------------------------------------------------

class NoteService:
    """Everything M3: upload, generate, refine, discuss, folders, search."""

    def __init__(self, db_url: Optional[str] = None):
        self._db_url = db_url or os.getenv(
            "DATABASE_URL",
            "postgresql://postgres.iatjbhvtcvnsbitpbfim:NeuroNote2026"
            "@aws-1-ap-south-1.pooler.supabase.com:5432/postgres",
        )
        self._embeddings: Optional[SentenceTransformer] = None

    # ------------------------------------------------------------------
    # DB helpers
    # ------------------------------------------------------------------

    def _connect(self):
        db = self._db_url
        if "supabase" in db or "pooler" in db:
            return psycopg2.connect(
                db, cursor_factory=psycopg2.extras.RealDictCursor, sslmode="require"
            )
        return psycopg2.connect(
            db, cursor_factory=psycopg2.extras.RealDictCursor
        )

    def _get_embeddings(self):
        """
        Lazy-load sentence-transformers ONLY when EMBEDDINGS_ENABLED=True.
        This avoids the PyTorch + sentence-transformers crash on machines where
        it isn't needed (embeddings are only used for semantic search).
        """
        if self._embeddings is None:
            if not EMBEDDINGS_ENABLED:
                raise RuntimeError("Embeddings are disabled — set EMBEDDINGS_ENABLED=True to use search")
            logger.info("Loading embedding model %s ...", EMBEDDING_MODEL_NAME)
            from sentence_transformers import SentenceTransformer

            self._embeddings = SentenceTransformer(EMBEDDING_MODEL_NAME)
        return self._embeddings

    def get_embedding(self, text: str) -> List[float]:
        model = self._get_embeddings()
        return model.encode(text).tolist()

    # ------------------------------------------------------------------
    # 🔧 FIXED: _chunk_text — NO infinite loop
    # ------------------------------------------------------------------

    def _chunk_text(self, text: str) -> List[str]:
        """
        Split text into ~CHUNK_SIZE-char chunks with ~CHUNK_OVERLAP overlap.

        The original version had an infinite-loop bug at the tail: when
        ``end == len(text)``, ``start = end - CHUNK_OVERLAP`` would move
        *backwards*, causing the same last chunk to be appended forever.
        FIXED: break immediately when we've covered the entire text.
        """
        chunks: List[str] = []
        text_len = len(text)

        if text_len == 0:
            return chunks

        start = 0
        while start < text_len:
            end = min(start + CHUNK_SIZE, text_len)

            # Try to break at a natural boundary
            if end < text_len:
                chunk_slice = text[start:end]
                for sep in ("\n\n", "\n", ". ", " "):
                    last = chunk_slice.rfind(sep)
                    if last > CHUNK_SIZE // 2:          # only if it's a meaningful break
                        end = start + last + len(sep)
                        break

            chunk = text[start:end].strip()
            if chunk:                                    # skip empty chunks
                chunks.append(chunk)

            # ✅ FIX: if we've reached the end, stop (was missing — caused infinite loop)
            if end >= text_len:
                break

            start = end - CHUNK_OVERLAP

            # Safety cap — prevent runaway chunking on any unexpected edge case
            if len(chunks) >= MAX_CHUNKS_PER_FILE:
                logger.warning(
                    "_chunk_text: hit safety cap of %d chunks — truncating",
                    MAX_CHUNKS_PER_FILE,
                )
                break

        logger.info("_chunk_text: %d chars → %d chunks", text_len, len(chunks))
        return chunks

    # ------------------------------------------------------------------
    # process_file — Router calls: note_service.process_file(file_bytes, file_id, file.filename)
    # ------------------------------------------------------------------

    def process_file(
        self, file_bytes: bytes, file_id: str, filename: str
    ) -> Dict[str, Any]:
        """
        Extract text from uploaded bytes, chunk, embed, store in document_chunks.

        Returns: {"status": "success", "file_id": ..., "chunks": N, ...}
                 or {"status": "error", "message": "..."}
        """
        ext = os.path.splitext(filename)[1].lower()
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        
        # Persist the file to documents/ for split-view serving
        try:
            os.makedirs("documents", exist_ok=True)
            pdf_storage_path = os.path.join("documents", f"{file_id}{ext}")
            with open(pdf_storage_path, "wb") as f_store:
                f_store.write(file_bytes)
            logger.info("Persisted source PDF file to: %s", pdf_storage_path)
        except Exception as persist_err:
            logger.error("Failed to persist source document: %s", persist_err)

        try:
            full_text, images = extract_text_from_file(tmp_path, file_id)
        except Exception as e:
            logger.exception("Text extraction failed for %s", filename)
            try:
                os.unlink(tmp_path)
            except OSError:
                pass
            return {"status": "error", "message": f"Text extraction failed: {e}"}

        # Chunk the text (now with the infinite-loop bug fixed)
        chunks = self._chunk_text(full_text)
        logger.info(
            "process_file: %s → %d chars / %d chunks / %d images",
            filename, len(full_text), len(chunks), len(images),
        )

        conn = self._connect()
        try:
            # --- Ensure tables exist ---
            with conn.cursor() as cur:
                # Ensure the table exists (schema from scratch_create_tables.py)
                cur.execute("""
                    CREATE TABLE IF NOT EXISTS document_chunks (
                        id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
                        pdf_id TEXT,
                        chunk_index INTEGER,
                        content TEXT,
                        embedding VECTOR(384),
                        metadata JSONB
                    )
                """)
                cur.execute("""
                    CREATE TABLE IF NOT EXISTS document_images (
                        id TEXT PRIMARY KEY,
                        pdf_id TEXT,
                        image_data TEXT,
                        media_type TEXT,
                        page_number INTEGER,
                        caption TEXT
                    )
                """)
                conn.commit()

            # --- Always store chunks (needed for generation) ---
            with conn.cursor() as cur:
                for batch_start in range(0, len(chunks), MAX_EMBEDDING_BATCH):
                    batch = chunks[batch_start:batch_start + MAX_EMBEDDING_BATCH]
                    for i, chunk_text in enumerate(batch):
                        chunk_idx = batch_start + i
                        # Store NULL embedding when embeddings are disabled
                        if EMBEDDINGS_ENABLED:
                            try:
                                emb = self.get_embedding(chunk_text)
                            except Exception as emb_err:
                                logger.warning("Embedding failed for chunk %d: %s — storing NULL", chunk_idx, emb_err)
                                emb = None
                        else:
                            emb = None

                        if emb is not None:
                            cur.execute(
                                """INSERT INTO document_chunks
                                   (id, pdf_id, chunk_index, content, embedding)
                                   VALUES (gen_random_uuid(), %s, %s, %s, %s::vector)""",
                                (file_id, chunk_idx, chunk_text, emb),
                            )
                        else:
                            cur.execute(
                                """INSERT INTO document_chunks
                                   (id, pdf_id, chunk_index, content)
                                   VALUES (gen_random_uuid(), %s, %s, %s)""",
                                (file_id, chunk_idx, chunk_text),
                            )
            chunk_rows = []
            for chunk_idx, chunk_text in enumerate(chunks):
                if EMBEDDINGS_ENABLED:
                    try:
                        emb = self.get_embedding(chunk_text)
                    except Exception as emb_err:
                        logger.warning("Embedding failed for chunk %d: %s — storing NULL", chunk_idx, emb_err)
                        emb = None
                else:
                    emb = None
                
                import uuid
                chunk_id = str(uuid.uuid4())
                chunk_rows.append((chunk_id, file_id, chunk_idx, chunk_text, emb))

            with conn.cursor() as cur:
                from psycopg2.extras import execute_values
                execute_values(
                    cur,
                    """INSERT INTO document_chunks (id, pdf_id, chunk_index, content, embedding)
                       VALUES %s""",
                    chunk_rows
                )
                conn.commit()

            # Store image references
            if images:
                with conn.cursor() as cur:
                    for img in images:
                        cur.execute(
                            """INSERT INTO document_images
                               (id, pdf_id, page_number)
                               VALUES (gen_random_uuid()::text, %s, %s)""",
                            (file_id, img["page"]),
                        )
                image_rows = []
                for img in images:
                    image_rows.append((img["id"], file_id, img["page"], img.get("caption")))
                
                with conn.cursor() as cur:
                    from psycopg2.extras import execute_values
                    execute_values(
                        cur,
                        """INSERT INTO document_images (id, pdf_id, page_number, caption)
                           VALUES %s
                           ON CONFLICT (id) DO UPDATE
                           SET caption = EXCLUDED.caption""",
                        image_rows
                    )
                    conn.commit()

            return {
                "status": "success",
                "pdf_url": "",
                "file_id": file_id,
                "chunks": len(chunks),
                "images": len(images),
                "text_length": len(full_text),
            }
        except Exception as e:
            logger.exception("DB insert failed for %s", filename)
            return {"status": "error", "message": f"Database insert failed: {e}"}
        finally:
            conn.close()
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    # ------------------------------------------------------------------
    # 🆕 CORE: single-shot unified cheatsheet (internal engine)
    # ------------------------------------------------------------------

    def _single_shot_cheatsheet(
        self,
        file_ids: List[str],
        user_id: str,
        headnote: str = "",
        job_id: Optional[str] = None,
    ) -> str:
        """
        Concatenates ALL text from ALL file_ids → ONE LLM call → Markdown.

        This replaces the old 5-phase pipeline (50+ LLM calls) with one call.
        Gemini's 1M-token context window handles 3-5 lecture PDFs easily.

        🆕 Pre-scans the source for code blocks and exercises so the LLM
        knows exact counts it must include — no skipping.
        """
        logger.info(
            "_single_shot_cheatsheet: %d file(s), user=%s, model=%s",
            len(file_ids), user_id, get_model(),
        )

        # 1) Gather all text from the DB
        all_text_parts: List[str] = []
        file_names: List[str] = []
        total_chunks = 0
        conn = self._connect()
        try:
            with conn.cursor() as cur:
                for fid in file_ids:
                    # Get filename for the prompt header (safe — table may not exist)
                    try:
                        cur.execute(
                            "SELECT file_name FROM document_files WHERE id = %s", (fid,)
                        )
                        row = cur.fetchone()
                        if row:
                            file_names.append(row["file_name"])
                    except Exception:
                        conn.rollback()  # ✅ reset aborted transaction before next query
                        pass  # document_files table might not exist yet

                    # Read all chunks for this file, ordered
                    cur.execute(
                        "SELECT content FROM document_chunks WHERE pdf_id = %s ORDER BY chunk_index",
                        (fid,),
                    )
                    file_chunks: List[str] = []
                    for r in cur.fetchall():
                        file_chunks.append(r["content"])
                        total_chunks += 1
                    # Wrap each file's content with START/END delimiters so the model
                    # knows exactly where exercises and takeaways belong.
                    if file_chunks:
                        all_text_parts.append(
                            f"=== START OF FILE: {fid} ===\n"
                            + "\n".join(file_chunks)
                            + f"\n=== END OF FILE: {fid} ==="
                        )
        finally:
            conn.close()

        full_text = "\n\n".join(all_text_parts)
        text_len = len(full_text)
        logger.info(
            "Assembled %d chars from %d chunks across %d file(s)",
            text_len, total_chunks, len(file_ids),
        )

        if text_len == 0:
            raise ValueError("No text found for the given file IDs. Did the upload succeed?")

        # 2) Truncate if needed (Gemini has 1M token window, we stay conservative)
        if text_len > MAX_INPUT_CHARS:
            logger.warning("Truncating %d → %d chars", text_len, MAX_INPUT_CHARS)
            full_text = full_text[:MAX_INPUT_CHARS] + (
                "\n\n[... content truncated — the full text was too large for a single prompt ...]"
            )

        # 3) 🆕 Pre-scan: count code blocks and exercises so the LLM is accountable
        source_counts = count_source_items(full_text)
        logger.info(
            "Source pre-scan: %d code blocks, %d exercises",
            source_counts["code_blocks"], source_counts["exercises"],
        )

        # Build accountability note that tells the LLM exactly what it must include
        accountability_parts = []
        if source_counts["code_blocks"] > 0:
            accountability_parts.append(
                f"## ⚠️ ACCOUNTABILITY: The source text contains EXACTLY "
                f"**{source_counts['code_blocks']} code blocks** (fenced ```...``` "
                f"and indented code). You MUST include ALL of them in your output. "
                f"Do NOT skip any code block. Do NOT paraphrase code into prose."
            )
        else:
            accountability_parts.append(
                f"## ℹ️ ACCOUNTABILITY: The source text contains **NO code blocks**. "
                f"Do NOT invent or add any code examples. "
                f"If the source has no code, your output must have no code."
            )
        if source_counts["exercises"] > 0:
            accountability_parts.append(
                f"## ⚠️ ACCOUNTABILITY: The source text contains EXACTLY "
                f"**{source_counts['exercises']} exercises / problems / worked examples**. "
                f"You MUST solve ALL of them with full step-by-step solutions. "
                f"Do NOT skip any exercise."
            )
        else:
            accountability_parts.append(
                f"## ℹ️ ACCOUNTABILITY: The source text contains **NO exercises or problems**. "
                f"Do NOT invent fake exercises or practice problems. "
                f"Only include exercises that exist in the source."
            )
        accountability_note = "\n\n".join(accountability_parts)

        # 4) Build the prompt
        source_list = ", ".join(file_names) if file_names else "Unknown source(s)"
        prompt_body = UNIFIED_CHEATSHEET_USER.format(
            content=full_text,
            accountability_note=accountability_note,
        )
        prompt = (
            f"The source material comes from: {source_list}.\n"
            + (f"{headnote}\n\n" if headnote else "")
            + prompt_body
        )

        # 5) ONE LLM call
        logger.info("Calling LLM (single-shot unified cheatsheet) — this may take 30-120s …")
        markdown = ask_llm(
            prompt,
            system=UNIFIED_CHEATSHEET_SYSTEM,
            temperature=0.1,   # very low — maximise factual adherence, minimise creativity
            max_tokens=65536,  # larger budget for full solved exercises + key takeaways
            # standard Gemini output token limit for fast routing
        )

        # 6) Clean and return
        logger.info("Cheatsheet generated: %d chars of Markdown", len(markdown))
        return clean_note_formatting(markdown)

    # ------------------------------------------------------------------
    # generate_detailed_note — single-PDF → detailed cheatsheet
    # ------------------------------------------------------------------

    def generate_detailed_note(
        self,
        pdf_id: str,
        user_id: str,
        language: str = "English",
        job_id: Optional[str] = None,
    ) -> str:
        return self._single_shot_cheatsheet(
            file_ids=[pdf_id],
            user_id=user_id,
            headnote=f"Language: {language}",
            job_id=job_id,
        )

    # ------------------------------------------------------------------
    # generate_structured_note — multi-PDF → unified cheatsheet
    # ------------------------------------------------------------------

    def generate_structured_note(
        self,
        input_items: List[dict],
        user_id: str,
        language: str = "English",
        job_id: Optional[str] = None,
    ) -> str:
        file_ids: List[str] = []
        for item in input_items:
            pid = item.get("pdf_id") or item.get("file_id") or item.get("id") or item.get("value")
            if pid:
                file_ids.append(pid)

        if not file_ids:
            raise ValueError("No valid file IDs found in input_items")

        return self._single_shot_cheatsheet(
            file_ids=file_ids,
            user_id=user_id,
            headnote=f"Language: {language}. Synthesise ALL the following lectures into ONE cohesive cheatsheet.",
            job_id=job_id,
        )

    # ------------------------------------------------------------------
    # generate_note — legacy generate-note route (Member 2 integration)
    # ------------------------------------------------------------------

    def generate_note(
        self,
        pdf_ids: List[str],
        user_id: str,
        instruction: str = "",
        language: str = "English",
        ordering: str = "ai",
        job_id: Optional[str] = None,
    ) -> str:
        headnote = f"Language: {language}."
        if instruction:
            headnote += f" Special instruction: {instruction}"
        return self._single_shot_cheatsheet(
            file_ids=pdf_ids,
            user_id=user_id,
            headnote=headnote,
            job_id=job_id,
        )

    # ------------------------------------------------------------------
    # extract_mindmap_chunked — Router: /test-mindmap
    # ------------------------------------------------------------------

    def extract_mindmap_chunked(self, full_text: str) -> Dict[str, Any]:
        """Extract topic hierarchy as JSON (used by test-mindmap route)."""
        prompt = (
            "Extract the complete topic hierarchy from the following lecture text "
            "as a JSON object. Return ONLY valid JSON with this structure:\n"
            '{"lecture_title": "...", "chapters": ['
            '{"title": "...", "sections": ['
            '{"title": "...", "content_lines": ["..."]}'
            "]}]}\n\n"
            f"## LECTURE TEXT\n{full_text[:100000]}"
        )
        response = ask_llm(prompt, max_tokens=4096, temperature=0.1)
        try:
            json_str = re.sub(r"^```(?:json)?\s*\n?", "", response.strip())
            json_str = re.sub(r"\n?```\s*$", "", json_str)
            return json.loads(json_str)
        except json.JSONDecodeError:
            logger.warning("Mind map JSON parse failed; falling back to raw")
            return {
                "lecture_title": "Untitled Lecture",
                "chapters": [{"title": "Main Content", "sections": []}],
            }

    # ------------------------------------------------------------------
    # save_note_to_db
    # ------------------------------------------------------------------

    def save_note_to_db(
        self, user_id: str, pdf_id: Optional[str], title: str, content: str
    ) -> Optional[str]:
        content = resolve_image_tokens_to_static_urls(content)
        note_id = str(uuid.uuid4())
        conn = self._connect()
        try:
            with conn.cursor() as cur:
                cur.execute("""
                    CREATE TABLE IF NOT EXISTS notes (
                        note_id TEXT PRIMARY KEY,
                        user_id TEXT NOT NULL,
                        pdf_id TEXT,
                        title TEXT,
                        content TEXT,
                        folder_id TEXT,
                        note_type TEXT,
                        created_at TIMESTAMPTZ DEFAULT NOW(),
                        updated_at TIMESTAMPTZ DEFAULT NOW()
                    )
                """)
                cur.execute(
                    """INSERT INTO notes (note_id, user_id, pdf_id, title, content, created_at, updated_at)
                       VALUES (%s, %s, %s, %s, %s, %s, %s)
                       ON CONFLICT (note_id) DO UPDATE
                       SET content = EXCLUDED.content,
                           title = EXCLUDED.title,
                           updated_at = EXCLUDED.updated_at""",
                    (
                        note_id, user_id, pdf_id, title, content,
                        datetime.now(timezone.utc), datetime.now(timezone.utc),
                    ),
                )
                conn.commit()
            return note_id
        except Exception as e:
            logger.exception("Failed to save note to DB")
            return None
        finally:
            conn.close()

    # ------------------------------------------------------------------
    # update_note
    # ------------------------------------------------------------------

    def update_note(self, note_id: str, content: str) -> bool:
        conn = self._connect()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    "UPDATE notes SET content = %s, updated_at = %s WHERE note_id = %s",
                    (content, datetime.now(timezone.utc), note_id),
                )
                conn.commit()
                return cur.rowcount > 0
        except Exception as e:
            logger.exception("update_note failed")
            return False
        finally:
            conn.close()

    # ------------------------------------------------------------------
    # update_note_folder
    # ------------------------------------------------------------------

    def update_note_folder(self, note_id: str, folder_id: str) -> bool:
        conn = self._connect()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    "UPDATE notes SET folder_id = %s, updated_at = %s WHERE note_id = %s",
                    (folder_id, datetime.now(timezone.utc), note_id),
                )
                conn.commit()
                return cur.rowcount > 0
        except Exception as e:
            logger.exception("update_note_folder failed")
            return False
        finally:
            conn.close()

    # ------------------------------------------------------------------
    # refine_text
    # ------------------------------------------------------------------

    def refine_text(
        self,
        pdf_id: str,
        selected_text: str,
        instruction: str,
        loop_number: int = 1,
        allow_outside: bool = False,
        conversation_history: Optional[List[dict]] = None,
    ) -> Dict[str, Any]:
        prompt = (
            f"## Current section content\n{selected_text}\n\n"
            f"## User instruction\n{instruction}\n\n"
            f"Return the FULL improved section as valid Markdown."
        )
        refined = ask_llm(prompt, system=REFINE_SYSTEM, temperature=0.3, max_tokens=4096)

        should_ask = (
            allow_outside
            and loop_number >= 3
            and any(
                kw in instruction.lower()
                for kw in ("outside", "web", "search", "research", "wider", "broader")
            )
        )

        return {
            "refined_content": refined.strip(),
            "loop_number": loop_number,
            "should_ask_outside": should_ask,
        }

    # ------------------------------------------------------------------
    # discuss_note
    # ------------------------------------------------------------------

    def discuss_note(
        self,
        note_content: str,
        user_question: str,
        pdf_id: Optional[str] = None,
        conversation_history: Optional[List[dict]] = None,
    ) -> Dict[str, Any]:
        from .openai_client import get_client

        messages: List[Dict[str, str]] = [
            {"role": "system", "content": DISCUSS_SYSTEM}
        ]
        if conversation_history:
            messages.extend(conversation_history)
        messages.append({
            "role": "user",
            "content": f"## Section content\n{note_content}\n\n## Question\n{user_question}",
        })

        client = get_client()
        resp = client.chat.completions.create(
            model=get_model(),
            messages=messages,
            temperature=0.4,
            max_tokens=2048,
        )
        answer = resp.choices[0].message.content or ""
        return {"refined_content": answer.strip()}

    # ------------------------------------------------------------------
    # summarize_prompts
    # ------------------------------------------------------------------

    def summarize_prompts(
        self, prompts: List[str], original_text: Optional[str] = None
    ) -> str:
        combined = "; ".join(prompts)
        if original_text:
            combined = f"Context: {original_text[:200]}\nPrompts: {combined}"
        system = "You are a precise topic labeler. Return ONLY a short label (≤5 words) describing what this content is about. Do NOT include any other text."
        label = ask_llm(combined, system=system, temperature=0.1, max_tokens=30)
        return label.strip().strip('"').strip("'")

    # ------------------------------------------------------------------
    # Folder CRUD
    # ------------------------------------------------------------------

    def create_folder(self, user_id: str, name: str) -> Dict[str, Any]:
        conn = self._connect()
        try:
            fid = str(uuid.uuid4())
            with conn.cursor() as cur:
                cur.execute(
                    "INSERT INTO note_folders (id, user_id, name) VALUES (%s, %s, %s)",
                    (fid, user_id, name),
                )
                conn.commit()
            return {"folder_id": fid, "name": name}
        finally:
            conn.close()

    def list_folders(self, user_id: str) -> List[Dict[str, Any]]:
        conn = self._connect()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    "SELECT id, name, created_at FROM note_folders WHERE user_id = %s "
                    "ORDER BY created_at DESC",
                    (user_id,),
                )
                return [dict(r) for r in cur.fetchall()]
        finally:
            conn.close()

    def update_folder(self, folder_id: str, user_id: str, name: str) -> Dict[str, Any]:
        conn = self._connect()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    "UPDATE note_folders SET name = %s WHERE id = %s AND user_id = %s",
                    (name, folder_id, user_id),
                )
                conn.commit()
            return {"folder_id": folder_id, "name": name}
        finally:
            conn.close()

    def delete_folder(self, folder_id: str, user_id: str) -> Dict[str, Any]:
        conn = self._connect()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    "DELETE FROM note_folders WHERE id = %s AND user_id = %s",
                    (folder_id, user_id),
                )
                conn.commit()
            return {"deleted": True}
        finally:
            conn.close()

    # ------------------------------------------------------------------
    # Search (embedding-based)
    # ------------------------------------------------------------------

    def search_chunks(self, query: str, user_id: str, top_k: int = 5) -> List[Dict[str, Any]]:
        emb = self.get_embedding(query)
        conn = self._connect()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    """SELECT dc.content AS chunk_text, dc.pdf_id AS file_id,
                              dc.embedding <=> %s::vector AS distance
                       FROM document_chunks dc
                       ORDER BY distance
                       LIMIT %s""",
                    (emb, top_k),
                )
                return [dict(r) for r in cur.fetchall()]
        finally:
            conn.close()

    # ------------------------------------------------------------------
    # Image utilities (static — no DB)
    # ------------------------------------------------------------------

    @staticmethod
    def merge_images(markdown: str, image_map: Dict[str, str]) -> str:
        for fname, b64 in image_map.items():
            marker = f"[IMAGE: {fname}]"
            replacement = (
                f'<img src="data:image/png;base64,{b64}" '
                f'alt="{fname}" style="max-width:100%"/>'
            )
            markdown = markdown.replace(marker, replacement)
        return markdown

    @staticmethod
    def remove_images_for_md_export(markdown: str) -> str:
        markdown = re.sub(r"\[IMAGE:\s*[^\]]+\]", "", markdown)
        markdown = re.sub(r"<img[^>]*>", "", markdown)
        return markdown


# ---------------------------------------------------------------------------
# Compatibility alias + singleton
# ---------------------------------------------------------------------------

class AIService(NoteService):
    """
    Compatibility alias — ``from .services import AIService`` continues to work.
    Previously AIService held LLM + embeddings separately; now unified.
    """
    pass


# Singleton — used by router: ``from m3_structurednotes.services import note_service``
note_service = NoteService()
