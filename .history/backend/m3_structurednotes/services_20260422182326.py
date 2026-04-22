import os
import json
from io import BytesIO
# Make specific dependencies optional
# They will be imported inside functions so other team members
# don't need to install them if they don't do AI logic.
from dotenv import load_dotenv
from .database import get_db_connection
from psycopg2.extras import RealDictCursor
import uuid
import math
import re
from collections import Counter
import json
import numpy as np

load_dotenv()

class AIService:

    def __init__(self):
        # Use HuggingFace local embeddings (free)
        self._embeddings = None
        # Use Groq for LLM (free tier)
        self._llm = None

    @property
    def embeddings(self):
        if self._embeddings is None:
            # OpenRouter does not provide free robust text-embedding models reliably.
            # We revert to local HuggingFace to keep document processing functional and free.
            from langchain_huggingface import HuggingFaceEmbeddings
            print("Loading HuggingFace Embeddings for Document Vectorization...")
            self._embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
        return self._embeddings

    @property
    def llm(self):
        if self._llm is None:
            from langchain_openai import ChatOpenAI
            
            # Use OpenRouter - dynamically read from environment
            api_key = os.getenv("OPENROUTER_API_KEY")
            model_name = os.getenv("OPENROUTER_MODEL", "liquid/lfm-2.5-1.2b-instruct:free")
            
            print(f"[*] Initializing AIService with model: {model_name}")
            
            self._llm = ChatOpenAI(
                api_key=api_key, 
                base_url="https://openrouter.ai/api/v1",
                model=model_name,
                temperature=0
            )
        return self._llm

    def extract_keywords_tfidf(self, text, top_n=10):
        """
        MANUAL ALGORITHM: Deep Topic Modeling via TF-IDF Scoring.
        We avoided using external libraries (like scikit-learn) to implement this completely manually.
        This algorithm calculates Term Frequency (TF) and Inverse Document Frequency (IDF) mathematically.
        By analyzing term rarity, it dives deeper than simple summarizing algorithms by pulling out core,
        sometimes hidden, thematic keywords (not just common words) from user content to build a highly structured note.
        """
        # Stopwords to ignore
        stopwords = {"the", "is", "in", "and", "to", "of", "a", "for", "on", "with", "as", "by", "this", 
                     "that", "it", "at", "from", "or", "an", "be", "are", "can", "will", "which"}
        
        # 1. Clean and tokenize text into words
        words = re.findall(r'\b[a-z]{3,}\b', text.lower())
        words = [w for w in words if w not in stopwords]
        
        if not words:
            return []
            
        # 2. Treat paragraphs/sentences as 'documents' for IDF calculation
        sentences = [s.strip() for s in text.lower().split('.') if len(s.strip()) > 10]
        if not sentences:
            sentences = [text.lower()]
            
        N = len(sentences)
        
        # 3. Calculate Term Frequency (TF) for the entire text
        word_counts = Counter(words)
        total_words = len(words)
        
        tfidf_scores = {}
        for word, count in word_counts.items():
            tf = count / total_words
            
            # 4. Calculate Inverse Document Frequency (IDF)
            # Count how many sentences contain the word
            doc_freq = sum(1 for sentence in sentences if word in sentence)
            idf = math.log(N / (1 + doc_freq)) + 1  # Add 1 to avoid zero weights
            
            # 5. Final TF-IDF Score
            tfidf_scores[word] = tf * idf
            
        # Sort words by score descending
        sorted_keywords = sorted(tfidf_scores.items(), key=lambda item: item[1], reverse=True)
        return [word for word, score in sorted_keywords[:top_n]]

    def process_pdf(self, file_bytes, pdf_id, original_filename):
        # Determine file type
        is_pptx = original_filename.lower().endswith('.pptx')
        extension = ".pptx" if is_pptx else ".pdf"

        # 1. Save File to disk for viewing
        save_dir = "documents"
        os.makedirs(save_dir, exist_ok=True)
        
        safe_filename = f"{pdf_id}{extension}"
        file_path = os.path.join(save_dir, safe_filename)
        
        with open(file_path, "wb") as f:
            f.write(file_bytes)
            
        # 2. Extract text with optional dependencies
        full_text = ""
        try:
            if is_pptx:
                from pptx import Presentation
                prs = Presentation(BytesIO(file_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            full_text += shape.text + " "
                    full_text += "\n"
            elif original_filename.lower().endswith(('.md', '.txt')):
                full_text = file_bytes.decode('utf-8', errors='ignore')
            else:
                from PyPDF2 import PdfReader
                reader = PdfReader(BytesIO(file_bytes))
                for page in reader.pages:
                    full_text += page.extract_text() or ""

            # 2b. Extract Images as Contextual Assets (PDF only)
            extracted_images = []
            if not is_pptx:
                try:
                    import fitz  # PyMuPDF
                    images_dir = "images"
                    os.makedirs(images_dir, exist_ok=True)
                    doc = fitz.open(stream=file_bytes, filetype="pdf")
                    img_count = 0
                    for page_num, page in enumerate(doc):
                        image_list = page.get_images(full=True)
                        for img_index, img in enumerate(image_list):
                            if img_count >= 10:  # max 10 images per doc
                                break
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            image_bytes = base_image["image"]
                            image_ext = base_image["ext"]
                            # Only save meaningful images (skip tiny icons < 5KB)
                            if len(image_bytes) < 5000:
                                continue
                            img_filename = f"{pdf_id}_p{page_num}_i{img_index}.{image_ext}"
                            img_path = os.path.join(images_dir, img_filename)
                            with open(img_path, "wb") as img_file:
                                img_file.write(image_bytes)
                            extracted_images.append({
                                "filename": img_filename,
                                "page": page_num + 1,
                                "url": f"/images/{img_filename}"
                            })
                            img_count += 1
                    doc.close()
                    if extracted_images:
                        print(f"Extracted {len(extracted_images)} contextual images from {original_filename}")
                except ImportError:
                    print("PyMuPDF not installed - skipping image extraction. Run: pip install pymupdf")
                except Exception as img_err:
                    print(f"Image extraction warning: {img_err}")

            # 3. Chunk text for Vector Storage
            from langchain_text_splitters import RecursiveCharacterTextSplitter
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len
            )
            chunks = text_splitter.split_text(full_text)
        except ImportError as e:
            return {
                "status": "error",
                "error": f"Missing optional AI dependency: {e}"
            }
        
        # 4. Generate Embeddings and Save to Postgres (pgvector)
        conn = None
        try:
            conn = get_db_connection()
            if conn:
                cur = conn.cursor()
                
                # Check if document_chunks table exists (lightweight check/fallback)
                # In production, schema migration handles this.
                
                # Generate embeddings in batch
                embeddings = self.embeddings.embed_documents(chunks)
                
                for i, (chunk_text, embedding) in enumerate(zip(chunks, embeddings)):
                    chunk_id = str(uuid.uuid4())
                    # Convert to pgvector format: "[0.1, 0.2, ...]"
                    embedding_str = "[" + ",".join(str(x) for x in embedding) + "]"
                    cur.execute("""
                        INSERT INTO document_chunks (id, pdf_id, chunk_index, content, embedding, metadata)
                        VALUES (%s, %s, %s, %s, %s::vector, %s)
                    """, (
                        chunk_id, 
                        pdf_id, 
                        i, 
                        chunk_text, 
                        embedding_str,  # proper pgvector format
                        json.dumps({"source": original_filename})
                    ))
                
                conn.commit()
                cur.close()
                conn.close()
                print(f"Stored {len(chunks)} chunks in Postgres (pgvector) for {pdf_id}")
            else:
                print("DB Connection failed, skipping vector storage.")
                
        except Exception as e:
            print(f"Error storing vectors in Postgres: {e}")
            if conn:
                conn.rollback()
        
        # 5. Return result
        return {
            "status": "success",
            "pdf_url": f"/documents/{safe_filename}",
            "extracted_text": full_text,
            "extracted_images": extracted_images  # Contextual Assets for AI embedding
        }

    def save_note_to_db(self, user_id, pdf_id, title, content):
        conn = get_db_connection()
        if not conn:
            print("Database connection failed")
            return None
        
        try:
            cur = conn.cursor()
            note_id = str(uuid.uuid4())
            
            # Use pdf_id if provided (ensure column exists or handle gracefully)
            cur.execute("""
                INSERT INTO notes (note_id, user_id, title, content, created_at, updated_at, note_type, is_in_folder, has_embeddings)
                VALUES (%s, %s, %s, %s, NOW(), NOW(), 'ai_generated', FALSE, TRUE)
                RETURNING note_id
            """, (note_id, user_id, title, content))
            
            returned_id = cur.fetchone()[0]
            
            conn.commit()
            cur.close()
            conn.close()
            print(f"Note saved to DB via RETURNING: {returned_id}")
            return returned_id
        except Exception as e:
            print(f"Error saving to DB: {e}")
            if conn:
                conn.rollback()
            return None

        
    def generate_note(self, pdf_id, user_id, instruction=None, full_text=None, language="English", extracted_images=None):
        # 1. Manual Algorithm: Extract Keywords from the full text
        # If full text isn't passed, we'd normally pull it from DB, but for speed we assume it's passed or extract simple context
        
        # Retrieve ALL relevant chunks to maintain the lecture flow (not just top 5)
        # We want a comprehensive note, so we fetch more chunks
        context_text = ""
        try:
            conn = get_db_connection()
            if conn:
                cur = conn.cursor()
                # Fetch chunks ordered by their original index to maintain lecture flow!
                cur.execute("""
                    SELECT content 
                    FROM document_chunks 
                    WHERE pdf_id = %s 
                    ORDER BY chunk_index ASC
                """, (pdf_id,))
                
                results = cur.fetchall()
                context_chunks = [row[0] for row in results]
                
                # --- NEW: Semantic Synthesizer Outline Builder ---
                # Manual algorithm: Maps chronological memory chunks into distinct "Knowledge Modules".
                # This explicitly forces the AI synthesis step to combine different subjects 
                # (like two uploaded notebooks) and draw novel, deeper ideas between structural boundaries!
                context_text = ""
                module_counter = 1
                for i, chunk in enumerate(context_chunks):
                    # Group every 3 chunks into a deeper thematic Module
                    if i % 3 == 0:  
                        context_text += f"\n\n=== [KNOWLEDGE MODULE {module_counter} - Focus on Depth & Cross-Linking Novel Ideas] ===\n"
                        module_counter += 1
                    context_text += f"{chunk}\n"
                
                # For presentations: keep ALL content (no truncation)
                # Allow up to 50000 chars to include full 20+ slide presentations
                MAX_API_CHARS = 50000
                if len(context_text) > MAX_API_CHARS:
                    print(f"Note: Presentation is very large ({len(context_text)} chars). Processing all available content...")
                    # Still include all content - the AI model can handle it
                else:
                    print(f"Context prepared: {len(context_text)} characters from full presentation")
                    
                cur.close()
                conn.close()
            else:
                # [OFFLINE FALLBACK]: Temporarily bypass database if not connected
                print(f"Bypassing DB: Manually parsing documents/{pdf_id}.pdf")
                from PyPDF2 import PdfReader
                try:
                    pdf = PdfReader(f"documents/{pdf_id}.pdf")
                    fallback_text = ""
                    for page in pdf.pages:
                        fallback_text += page.extract_text()
                        
                    # Build simple synthetic chunks
                    from langchain_text_splitters import RecursiveCharacterTextSplitter
                    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
                    context_chunks = splitter.split_text(fallback_text)
                    
                    context_text = ""
                    module_counter = 1
                    for i, chunk in enumerate(context_chunks):
                        if i % 3 == 0:
                            context_text += f"\n\n=== [KNOWLEDGE MODULE {module_counter}] ===\n"
                            module_counter += 1
                        context_text += f"{chunk}\n"
                        
                    if len(context_text) > 14000:
                        context_text = context_text[:14000] + "\n...[Content truncated]"
                except Exception as e:
                    return f"Offline Fallback Error: {str(e)}"
        except Exception as e:
            print(f"Error retrieving vectors: {e}")
            return f"Error exploring document: {str(e)}"

        if not context_text:
            return "No relevant content found in the document."
            
        # Run our manual TF-IDF Algorithm!
        top_keywords = self.extract_keywords_tfidf(context_text, top_n=10)
        keyword_str = ", ".join(top_keywords)
        print(f"Manual Algorithm Extracted Keywords: {keyword_str}")
            
        # Build image context for the AI
        image_instruction = ""
        if extracted_images:
            backend_base = "http://127.0.0.1:8000"
            image_lines = []
            for img in extracted_images:
                full_url = backend_base + img['url']
                image_lines.append(f"- Page {img['page']}: ![Contextual Image]({full_url})")
            image_context = "\n".join(image_lines)
            image_instruction = f"""
        5. CONTEXTUAL ASSETS (MANDATORY): The following images were extracted directly from the document.
           You MUST embed each image into the most semantically relevant Knowledge Module using this exact Markdown syntax:
           ![Descriptive Caption](image_url)
           Each image should appear ONCE, placed naturally inside the module where it best supports the text.
           Available images:
{image_context}
        """

        # Create a highly structured prompt for professional, clean notes
        prompt = f"""
You are an expert note-taking specialist creating a professional, high-quality study guide.

CRITICAL FORMATTING RULES (MANDATORY - FOLLOW EXACTLY):
1. Use proper Markdown hierarchy: # for main title, ## for sections, ### for subsections
2. Every section MUST start with ## (section heading)
3. Use **bold** for key terms and concepts, never use excessive formatting
4. Use bullet points (-) for lists, not numbers for clarity
5. Add 1-2 line breaks between sections for readability
6. Keep paragraphs concise (3-4 sentences max per paragraph)
7. Add a "Key Takeaways" section at the end with 3-5 bullet points
8. Do NOT use code blocks unless showing actual code or data
9. Embed images naturally within relevant sections using Markdown syntax

CONTENT STRUCTURE (In this exact order):
1. Title: # {keyword_str.split(",")[0].title()} - Comprehensive Study Guide
2. Executive Summary (brief 2-3 sentence overview)
3. Main Sections (based on document flow - use ## for each):
   - For each major topic, explain clearly and concisely
   - Bold the key concepts
   - Use 1-2 supporting examples
4. Conceptual Connections (explain how topics relate)
5. Visual Aids (include at least one Mermaid diagram if appropriate)
{image_instruction}
6. Key Takeaways (bullet list of 3-5 main points)

TONE & STYLE:
- Professional, clear, and educational
- Avoid redundancy and excessive details
- Use active voice
- Write for someone learning this topic for the first time

DOCUMENT CONTENT:
{context_text}

CRITICAL: Output ONLY the Markdown note. No explanations, no meta-text. Start directly with the # title."""
        
        try:
            response = self.llm.invoke(prompt)
            return response.content
        except Exception as e:
            return f"Error generating note: {str(e)}"



    def update_note(self, note_id: str, new_content: str):
        """
        Updates the content of an existing note in the database.
        """
        conn = get_db_connection()
        if not conn:
            return False
            
        try:
            cur = conn.cursor()
            cur.execute(
                "UPDATE notes SET content = %s, updated_at = NOW() WHERE note_id = %s",
                (new_content, note_id)
            )
            rows_updated = cur.rowcount
            conn.commit()
            cur.close()
            conn.close()
            
            if rows_updated == 0:
                print(f"Warning: Attempted to update non-existent note {note_id}")
                return False
                
            return True
        except Exception as e:
            print(f"Error updating note: {e}")
            return False

    def refine_text(self, pdf_id: str, selected_text: str, instruction: str):
        """
        Refines or explains a specific piece of text using the original PDF context.
        """
        # 1. Retrieve relevant context from the PDF
        query_text = f"{selected_text} {instruction}"
        context_text = ""
        try:
            conn = get_db_connection()
            if conn:
                cur = conn.cursor()
                try:
                    # Try semantic vector search first
                    query_embedding = self.embeddings.embed_query(query_text)
                    embedding_str = "[" + ",".join(str(x) for x in query_embedding) + "]"
                    cur.execute("""
                        SELECT content 
                        FROM document_chunks 
                        WHERE pdf_id = %s 
                        ORDER BY embedding <=> %s::vector 
                        LIMIT 3
                    """, (pdf_id, embedding_str))
                except Exception as vec_err:
                    print(f"Vector search failed, falling back to sequential chunks: {vec_err}")
                    # Fallback: just get the first 3 chunks in order
                    cur.execute("""
                        SELECT content 
                        FROM document_chunks 
                        WHERE pdf_id = %s 
                        ORDER BY chunk_index ASC
                        LIMIT 3
                    """, (pdf_id,))
                
                results = cur.fetchall()
                context_chunks = [row[0] for row in results]
                context_text = "\n\n".join(context_chunks)
                cur.close()
                conn.close()
        except Exception as e:
            print(f"Error retrieving context for refinement: {e}")
            context_text = ""  # Fallback to no context
        
        # 2. Prompt the LLM
        prompt = f"""
        You are an expert tutor. The user has selected a specific part of a note and asked for a modification or explanation.
        
        CONTEXT from the original document:
        {context_text}
        
        USER SELECTED TEXT: "{selected_text}"
        USER INSTRUCTION: "{instruction}"
        
        YOUR TASK:
        Provide the requested refinement, explanation, or example based strictly on the provided CONTEXT. 
        Do not make up information not present in the context.
        If the instruction is to "rewrite", rewrite the selected text using the context.
        If the instruction is to "give an example", provide a concrete example derived from the context.
        Return ONLY the refined text or explanation.
        """
        
        # 3. Generate Response
        try:
            response = self.llm.invoke(prompt)
            return {"refined_content": response.content}
        except Exception as e:
            return {"error": str(e)}

    def summarize_prompts(self, prompts: list, original_text: str = None) -> str:
        """
        Takes a chronological sequence of refinement instructions and generates a topic title
        plus a brief description of the overall idea behind the refinement journey.
        """
        if not prompts:
            return "Refined Section"
            
        prompt_history = " \u2794 ".join(prompts)
        
        orig_instruction = ""
        if original_text:
            cleaned_text = original_text[:120].replace('\n', ' ')
            orig_instruction = f"ORIGINAL TEXT SNIPPET (what the user was working on): '{cleaned_text}'"
            
        system_msg = f"""
        You are a study note assistant. A student refined a piece of text through multiple AI iterations.
        Analyze the journey and produce TWO things, exactly in this format:

        TITLE: <A short, clear 4-6 word title capturing the core topic (capitalize like a Title)>
        DESCRIPTION: <ONE sentence (max 20 words) explaining the overall idea or key insight explored across all the refinement steps>

        {orig_instruction}

        REFINEMENT SEQUENCE (what the student asked AI to do step by step):
        {prompt_history}

        Respond ONLY with the TITLE and DESCRIPTION lines. Nothing else.
        """
        
        try:
            response = self.llm.invoke(system_msg)
            raw = response.content.strip()
            # Parse the two lines
            title = "AI Refined Section"
            description = ""
            for line in raw.split('\n'):
                if line.strip().startswith('TITLE:'):
                    title = line.replace('TITLE:', '').strip().strip('"').strip("'")
                elif line.strip().startswith('DESCRIPTION:'):
                    description = line.replace('DESCRIPTION:', '').strip().strip('"').strip("'")
            # Return combined so frontend can split on '||'
            return f"{title}||{description}"
        except Exception as e:
            return "AI Refined Adjustment||"


    # --- Folder Management ---
    def create_folder(self, user_id, name):
        conn = get_db_connection()
        if not conn:
            return None
        
        try:
            cur = conn.cursor()
            folder_id = str(uuid.uuid4())
            cur.execute(
                "INSERT INTO folders (id, user_id, name) VALUES (%s, %s, %s) RETURNING id",
                (folder_id, user_id, name)
            )
            conn.commit()
            cur.close()
            conn.close()
            return {"id": folder_id, "name": name}
        except Exception as e:
            print(f"Error creating folder: {e}")
            return None

    def get_all_folders(self, user_id):
        conn = get_db_connection()
        if not conn:
            return []
        
        try:
            cur = conn.cursor(cursor_factory=RealDictCursor) # Ensure dicts are returned
            cur.execute("SELECT * FROM folders WHERE user_id = %s ORDER BY created_at DESC", (user_id,))
            folders = cur.fetchall()
            cur.close()
            conn.close()
            return folders
        except Exception as e:
            print(f"Error getting folders: {e}")
            return []

    def update_note_folder(self, note_id, folder_id):
        conn = get_db_connection()
        if not conn:
            return False
            
        try:
            cur = conn.cursor()
            cur.execute(
                "UPDATE notes SET folder_id = %s, updated_at = NOW() WHERE note_id = %s",
                (folder_id, note_id)
            )
            # Check if the note actually existed and was updated
            success = cur.rowcount > 0 
            
            conn.commit()
            cur.close()
            conn.close()
            return success
        except Exception as e:
            # Fixed the string formatting here (added the 'f' before the quotes)
            print(f"Error updating note folder: {e}") 
            if conn:
                conn.close()
            return False

    def get_all_notes(self, user_id, folder_id=None):
        conn = get_db_connection()
        if not conn:
            print("DB Connection failed in get_all_notes")
            return []
        
        try:
            print(f"Fetching notes for user: {user_id}, folder: {folder_id}")
            cur = conn.cursor(cursor_factory=RealDictCursor)
            
            query = """
                SELECT note_id, title, created_at, note_type, is_in_folder, folder_id 
                FROM notes 
                WHERE user_id = %s
            """
            params = [user_id]
            
            if folder_id:
                query += " AND folder_id = %s"
                params.append(folder_id)
                
            query += " ORDER BY updated_at DESC"
            
            cur.execute(query, tuple(params))
            notes = cur.fetchall()
            cur.close()
            conn.close()
            print(f"Found {len(notes)} notes")
            return notes
        except Exception as e:
            print(f"Error getting notes: {e}")
            if conn:
                conn.close()
            return []

ai_service = AIService()
