import hashlib
from io import BytesIO
from typing import List

import PyPDF2
import docx
from fastapi import HTTPException, UploadFile

from config.config import settings

# In-memory extraction cache keyed by SHA-256 of file bytes.
# Capped at _CACHE_MAX entries
from collections import OrderedDict as _OD
_CACHE_MAX = 100
_extraction_cache: _OD = _OD()


class FileProcessor:
    """
    Extracts plain text from all file types accepted by the frontend.

    Supported:
        Documents : pdf, doc, docx, txt, rtf
        Spreadsheets: xlsx, xls
        Presentations: ppt, pptx
        Images (OCR)  : jpg, jpeg, png, gif, webp, bmp, tiff
        E-books       : epub
    """

    # Every extension the frontend allows
    _ALL_SUPPORTED: set = {
        "pdf", "doc", "docx", "txt", "rtf",
        "xlsx", "xls",
        "ppt", "pptx",
        "jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff",
        "epub",
    }

    async def process_files(self, files: List[UploadFile]) -> str:
        """Process uploaded files and extract text, combining into one string.

        Results are cached by SHA-256 so repeated uploads of the same file
        skip extraction entirely.
        """
        combined_text = ""

        for file in files:
            ext = file.filename.split(".")[-1].lower()

            if ext not in self._ALL_SUPPORTED:
                print(f"⚠️  Skipping {file.filename} (unsupported type: .{ext})")
                continue

            raw = await file.read()
            cache_key = hashlib.sha256(raw).hexdigest()

            if cache_key in _extraction_cache:
                print(f"📋 Cache hit: {file.filename} ({len(raw):,} bytes)")
                text = _extraction_cache[cache_key]
            else:
                text = self._extract(raw, ext, file.filename)
                if text:
                    _extraction_cache[cache_key] = text
                    if len(_extraction_cache) > _CACHE_MAX:
                        _extraction_cache.popitem(last=False)  # evict oldest entry
                    print(f"📥 Cached extraction: {file.filename} ({len(text):,} chars)")

            if text:
                combined_text += f"\n\n--- {file.filename} ---\n{text}"
            else:
                print(f"ℹ️  No text extracted from {file.filename}")

        if not combined_text.strip():
            raise HTTPException(
                status_code=400,
                detail=(
                    "No text could be extracted from the uploaded files. "
                    "For images, ensure they contain readable text. "
                    "For documents, ensure they are not password-protected."
                ),
            )

        return combined_text.strip()

    # Router

    def _extract(self, content: bytes, ext: str, filename: str) -> str:
        """Route to the correct extractor based on file extension"""
        try:
            if ext == "pdf":
                return self._extract_from_pdf(content)
            elif ext in ("doc", "docx"):
                return self._extract_from_docx(content)
            elif ext == "txt":
                return content.decode("utf-8", errors="ignore")
            elif ext == "rtf":
                return self._extract_from_rtf(content)
            elif ext in ("xlsx", "xls"):
                return self._extract_from_excel(content)
            elif ext in ("ppt", "pptx"):
                return self._extract_from_pptx(content)
            elif ext in ("jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff"):
                return self._extract_from_image(content)
            elif ext == "epub":
                return self._extract_from_epub(content)
            else:
                return ""
        except Exception as e:
            print(f"⚠️  Extraction failed for {filename}: {e}")
            return ""

    # Extractors

    def _extract_from_pdf(self, content: bytes) -> str:
        """Extract text from a PDF file"""
        try:
            reader = PyPDF2.PdfReader(BytesIO(content))
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            # If PDF has no selectable text (scanned), fall back to OCR
            if not text.strip():
                print("ℹ️  PDF has no selectable text — attempting OCR")
                return self._ocr_pdf(content)
            return text
        except Exception as e:
            print(f"⚠️  PDF extraction failed: {e}")
            return ""

    def _ocr_pdf(self, content: bytes) -> str:
        """OCR a scanned PDF by rendering each page as an image"""
        try:
            import fitz  # PyMuPDF
            import pytesseract
            from PIL import Image

            doc = fitz.open(stream=content, filetype="pdf")
            text = ""
            for page in doc:
                pix = page.get_pixmap(dpi=200)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text += pytesseract.image_to_string(img) + "\n"
            return text
        except ImportError:
            print("⚠️  PyMuPDF or pytesseract not installed for scanned PDF OCR")
            return ""
        except Exception as e:
            print(f"⚠️  PDF OCR failed: {e}")
            return ""

    def _extract_from_docx(self, content: bytes) -> str:
        """Extract text from a DOCX file"""
        try:
            doc = docx.Document(BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            print(f"⚠️  DOCX extraction failed: {e}")
            return ""

    def _extract_from_rtf(self, content: bytes) -> str:
        """Extract text from an RTF file"""
        try:
            from striprtf.striprtf import rtf_to_text
            return rtf_to_text(content.decode("utf-8", errors="ignore"))
        except ImportError:
            print("⚠️  striprtf not installed. Run: pip install striprtf")
            # Fallback: strip RTF control words with regex
            import re
            raw = content.decode("utf-8", errors="ignore")
            text = re.sub(r"\{[^}]*\}", "", raw)
            text = re.sub(r"\\[a-z]+\d* ?", "", text)
            return text.strip()
        except Exception as e:
            print(f"⚠️  RTF extraction failed: {e}")
            return ""

    def _extract_from_excel(self, content: bytes) -> str:
        """Extract text from an XLSX/XLS file"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(BytesIO(content), read_only=True, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                text += f"\nSheet: {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) for cell in row if cell is not None)
                    if row_text:
                        text += row_text + "\n"
            return text
        except ImportError:
            print("⚠️  openpyxl not installed. Run: pip install openpyxl")
            return ""
        except Exception as e:
            print(f"⚠️  Excel extraction failed: {e}")
            return ""

    def _extract_from_pptx(self, content: bytes) -> str:
        """Extract text from a PPTX/PPT file"""
        try:
            from pptx import Presentation
            prs = Presentation(BytesIO(content))
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\nSlide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            return text
        except ImportError:
            print("⚠️  python-pptx not installed. Run: pip install python-pptx")
            return ""
        except Exception as e:
            print(f"⚠️  PPTX extraction failed: {e}")
            return ""

    def _extract_from_image(self, content: bytes) -> str:
        """Extract text from an image using Tesseract OCR"""
        try:
            import pytesseract
            from PIL import Image

            img = Image.open(BytesIO(content))

            # Convert to RGB if needed
            if img.mode not in ("RGB", "L"):
                img = img.convert("RGB")

            text = pytesseract.image_to_string(img, config="--psm 6")
            return text.strip()
        except ImportError:
            print(
                "⚠️  pytesseract or Pillow not installed.\n"
                "    Run: pip install pytesseract Pillow\n"
                "    Also install Tesseract OCR engine:\n"
                "      Ubuntu/Debian: sudo apt install tesseract-ocr\n"
                "      macOS:         brew install tesseract"
            )
            return ""
        except Exception as e:
            print(f"⚠️  Image OCR failed: {e}")
            return ""

    def _extract_from_epub(self, content: bytes) -> str:
        """Extract text from an EPUB file"""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            # write to a temp file as ebooklib requires file path 
            import os
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".epub", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name

            try:
                book = epub.read_epub(tmp_path)
                text = ""
                for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text += soup.get_text(separator="\n") + "\n"
                return text.strip()
            finally:
                os.unlink(tmp_path)

        except ImportError:
            print(
                "⚠️  ebooklib or beautifulsoup4 not installed.\n"
                "    Run: pip install ebooklib beautifulsoup4"
            )
            return ""
        except Exception as e:
            print(f"⚠️  EPUB extraction failed: {e}")
            return ""
