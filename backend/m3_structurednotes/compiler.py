import os
import re
import uuid
import logging
import tempfile
import numpy as np
import fitz  # PyMuPDF
from PIL import Image
from typing import List, Dict, Any, Optional

from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle, KeepTogether
)

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("structured_notes_compiler")

# ─────────────────────────────────────────────────────────────
# 1. CORE IMAGE UTILITIES (CROP, CHROME REMOVE, ASPECT RATIO)
# ─────────────────────────────────────────────────────────────

def get_content_bbox(image: Image.Image, bg_color: int = 255, tolerance: int = 15) -> Optional[tuple]:
    """
    Finds the bounding box of non-background content pixels in a grayscale image.
    bg_color: 255 for white background, 0 for black.
    Returns (x_min, y_min, x_max, y_max) or None if the image is empty background.
    """
    gray = image.convert("L")
    data = np.array(gray)
    
    if bg_color == 255:
        # Background is white, find pixels darker than threshold
        mask = data < (255 - tolerance)
    else:
        # Background is black, find pixels lighter than threshold
        mask = data > tolerance
        
    coords = np.argwhere(mask)
    if coords.size == 0:
        return None
        
    y_min, x_min = coords.min(axis=0)
    y_max, x_max = coords.max(axis=0)
    
    # Pad by a small margin
    padding = 4
    x_min = max(0, x_min - padding)
    y_min = max(0, y_min - padding)
    x_max = min(image.width, x_max + padding)
    y_max = min(image.height, y_max + padding)
    
    return (int(x_min), int(y_min), int(x_max), int(y_max))


def detect_repeating_chrome(images: List[Image.Image]) -> tuple:
    """
    Compares multiple page/slide images to find rows at the top or bottom that are identical.
    Identifies headers, footers, slide numbers, logos.
    Returns (top_ratio, bottom_ratio).
    """
    if len(images) < 2:
        return 0.0, 0.0
        
    h, w = 800, 800
    sample_images = images[:6]  # Analyze up to first 6 pages
    resized = [img.resize((w, h)).convert("L") for img in sample_images]
    
    stacked = np.stack([np.array(img) for img in resized], axis=0)  # shape (N, h, w)
    
    # Compute row-wise variance across pages
    row_variance = np.var(stacked, axis=0)  # shape (h, w)
    row_variance_mean = np.mean(row_variance, axis=1)  # shape (h,)
    
    # Identical or near-identical rows will have extremely low variance (e.g. < 4.0)
    is_static_row = row_variance_mean < 4.0
    
    # Check top 15% for headers
    top_trim = 0
    for i in range(int(h * 0.15)):
        if is_static_row[i]:
            top_trim = i + 1
        else:
            break
            
    # Check bottom 15% for footers / slide numbers
    bottom_trim = 0
    for i in range(h - 1, h - int(h * 0.15) - 1, -1):
        if is_static_row[i]:
            bottom_trim = h - i
        else:
            break
            
    return top_trim / h, bottom_trim / h


def clean_and_crop_image(
    image: Image.Image, top_chrome_ratio: float, bottom_chrome_ratio: float
) -> Image.Image:
    """
    Crops away detected top/bottom repeating chrome, then performs a content bounding box crop.
    """
    w, h = image.size
    
    # Crop repeating chrome first
    y_start = int(h * top_chrome_ratio)
    y_end = int(h * (1.0 - bottom_chrome_ratio))
    
    if y_end <= y_start:
        y_start, y_end = 0, h
        
    chrome_stripped = image.crop((0, y_start, w, y_end))
    
    # Find bounding box of remaining content
    bbox = get_content_bbox(chrome_stripped, bg_color=255, tolerance=15)
    if bbox:
        return chrome_stripped.crop(bbox)
        
class SourceInventory:
    def __init__(self):
        self.text_by_page: Dict[int, str] = {}
        self.images_by_page: Dict[int, List[str]] = {}
        self.all_text = ""
        self.file_id = str(uuid.uuid4())[:8]


def get_image_caption(img_path: str, page_text: str, skip_vision: bool) -> str:
    """
    Returns description/caption for the image. If skip_vision is True, uses a formatted placeholder.
    """
    filename = os.path.basename(img_path)
    base = os.path.splitext(filename)[0]
    
    if skip_vision:
        return f"Diagram extracted from source ({base})"
        
    try:
        from m3_structurednotes.services import describe_image
        with open(img_path, "rb") as f:
            image_bytes = f.read()
        caption = describe_image(image_bytes, page_text)
        return caption
    except Exception as e:
        logger.error("Vision descriptor failed: %s. Falling back to default caption.", e)
        return f"Diagram extracted from source ({base})"


def ingest_pdf(pdf_path: str, output_image_dir: str, skip_vision: bool = False) -> SourceInventory:
    """
    Ingests PDF, extracts text per page, and rasterizes pages to disk.
    """
    logger.info("Ingesting PDF: %s (skip_vision=%s)", pdf_path, skip_vision)
    inventory = SourceInventory()
    doc = fitz.open(pdf_path)
    os.makedirs(output_image_dir, exist_ok=True)
    
    page_images = []
    
    # Rasterize pages and extract text
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        text = page.get_text()
        inventory.text_by_page[page_idx + 1] = text
        
        # Render page to PNG at 150 DPI for visual analysis
        pix = page.get_pixmap(dpi=150)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        page_images.append(img)
        
    # Detect repeating chrome across pages
    top_chrome, bottom_chrome = detect_repeating_chrome(page_images)
    logger.info("Detected repeating chrome ratios: Top=%.2f, Bottom=%.2f", top_chrome, bottom_chrome)
    
    # Crop and save pages containing diagrams/diagram candidates
    for page_idx, img in enumerate(page_images):
        page_num = page_idx + 1
        text = inventory.text_by_page[page_num]
        cropped_img = clean_and_crop_image(img, top_chrome, bottom_chrome)
        
        # Save to temp location first to check descriptor
        img_filename = f"doc{inventory.file_id}_p{page_num}_img1.png"
        img_path = os.path.join(output_image_dir, img_filename)
        cropped_img.save(img_path)
        
        caption = get_image_caption(img_path, text, skip_vision)
        
        if caption == "SKIP_DECORATIVE_ICON":
            logger.info("Page %d image identified as decorative icon. Skipping.", page_num)
            try:
                os.remove(img_path)
            except OSError:
                pass
        else:
            token = f'[IMAGE:{img_filename[:-4]}|caption: "{caption}"]'
            inventory.text_by_page[page_num] += f"\n\n{token}\n"
            inventory.images_by_page[page_num] = [img_filename]
        
    doc.close()
    
    # Combine text
    inventory.all_text = "\n\n".join(
        f"--- Page {page_num} ---\n{text}"
        for page_num, text in inventory.text_by_page.items()
    )
    return inventory


def convert_pptx_to_pdf_libreoffice(pptx_path: str, output_dir: str) -> Optional[str]:
    """
    Attempts to run LibreOffice headless to convert PPTX to PDF.
    """
    import subprocess
    soffice_paths = [
        "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
    ]
    
    for path in soffice_paths:
        try:
            cmd = [path, "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path]
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
            pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
            if os.path.exists(pdf_path):
                return pdf_path
        except Exception:
            continue
    return None


def ingest_pptx(pptx_path: str, output_image_dir: str, skip_vision: bool = False) -> SourceInventory:
    """
    Ingests PPTX. Converts to PDF if LibreOffice is available, otherwise extracts shapes directly.
    """
    logger.info("Ingesting PPTX: %s (skip_vision=%s)", pptx_path, skip_vision)
    temp_dir = tempfile.mkdtemp()
    pdf_path = convert_pptx_to_pdf_libreoffice(pptx_path, temp_dir)
    
    if pdf_path:
        logger.info("Successfully converted PPTX to PDF using LibreOffice.")
        inventory = ingest_pdf(pdf_path, output_image_dir, skip_vision)
        # Clean up temp
        try:
            os.remove(pdf_path)
            os.rmdir(temp_dir)
        except OSError:
            pass
        return inventory
        
    # Headless fallback: extract pictures directly using python-pptx
    logger.warning("LibreOffice not found. Falling back to python-pptx direct extraction.")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    inventory = SourceInventory()
    prs = Presentation(pptx_path)
    os.makedirs(output_image_dir, exist_ok=True)
    
    slide_idx = 1
    for slide in prs.slides:
        slide_text_parts = []
        img_idx = 1
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_text_parts.append(paragraph.text.strip())
        slide_text = "\n".join(slide_text_parts)
        inventory.text_by_page[slide_idx] = slide_text
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                img_bytes = image.blob
                ext = image.ext or "png"
                
                img_filename = f"doc{inventory.file_id}_p{slide_idx}_img{img_idx}.{ext}"
                img_path = os.path.join(output_image_dir, img_filename)
                
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                    
                # Crop whitespace margins off the direct picture shape
                try:
                    pil_img = Image.open(img_path)
                    bbox = get_content_bbox(pil_img, bg_color=255, tolerance=15)
                    if bbox:
                        pil_img.crop(bbox).save(img_path)
                except Exception:
                    pass
                    
                caption = get_image_caption(img_path, slide_text, skip_vision)
                
                if caption == "SKIP_DECORATIVE_ICON":
                    try:
                        os.remove(img_path)
                    except OSError:
                        pass
                else:
                    token = f'[IMAGE:{img_filename[:-4]}|caption: "{caption}"]'
                    inventory.text_by_page[slide_idx] += f"\n\n{token}\n"
                    if slide_idx not in inventory.images_by_page:
                        inventory.images_by_page[slide_idx] = []
                    inventory.images_by_page[slide_idx].append(img_filename)
                    img_idx += 1
                
        slide_idx += 1
        
    inventory.all_text = "\n\n".join(
        f"--- Slide {slide_num} ---\n{text}"
        for slide_num, text in inventory.text_by_page.items()
    )
    return inventory


def ingest_docx(docx_path: str, output_image_dir: str, skip_vision: bool = False) -> SourceInventory:
    """
    Ingests DOCX. Extracts text paragraphs and embedded media shapes.
    """
    logger.info("Ingesting DOCX: %s (skip_vision=%s)", docx_path, skip_vision)
    from docx import Document
    inventory = SourceInventory()
    doc = Document(docx_path)
    os.makedirs(output_image_dir, exist_ok=True)
    
    # Extract text paragraphs
    text_parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    docx_text = "\n".join(text_parts)
    inventory.text_by_page[1] = docx_text
    
    # Extract embedded shapes
    img_idx = 1
    for rel_id, rel in doc.part.rels.items():
        if "image" in rel.target_ref:
            img_bytes = rel.target_part.blob
            ext = os.path.splitext(rel.target_ref)[1].replace(".", "") or "png"
            
            img_filename = f"doc{inventory.file_id}_img{img_idx}.{ext}"
            img_path = os.path.join(output_image_dir, img_filename)
            with open(img_path, "wb") as f:
                f.write(img_bytes)
                
            try:
                pil_img = Image.open(img_path)
                bbox = get_content_bbox(pil_img, bg_color=255, tolerance=15)
                if bbox:
                    pil_img.crop(bbox).save(img_path)
            except Exception:
                pass
                
            caption = get_image_caption(img_path, docx_text, skip_vision)
            
            if caption == "SKIP_DECORATIVE_ICON":
                try:
                    os.remove(img_path)
                except OSError:
                    pass
            else:
                token = f'[IMAGE:{img_filename[:-4]}|caption: "{caption}"]'
                inventory.text_by_page[1] += f"\n\n{token}\n"
                if 1 not in inventory.images_by_page:
                    inventory.images_by_page[1] = []
                inventory.images_by_page[1].append(img_filename)
                img_idx += 1
                
    inventory.all_text = inventory.text_by_page[1]
    return inventory


def ingest_folder(folder_path: str, output_image_dir: str, skip_vision: bool = False) -> SourceInventory:
    """
    Ingests a folder containing text files (.txt, .md) and images (.png, .jpg, .jpeg).
    """
    logger.info("Ingesting folder: %s (skip_vision=%s)", folder_path, skip_vision)
    import shutil
    inventory = SourceInventory()
    os.makedirs(output_image_dir, exist_ok=True)
    
    text_files = []
    image_files = []
    
    for filename in sorted(os.listdir(folder_path)):
        full_path = os.path.join(folder_path, filename)
        if os.path.isdir(full_path):
            continue
            
        ext = os.path.splitext(filename)[1].lower()
        if ext in (".txt", ".md"):
            text_files.append(full_path)
        elif ext in (".png", ".jpg", ".jpeg"):
            image_files.append(full_path)
            
    # Read and merge all texts
    merged_text_parts = []
    for tf in text_files:
        try:
            with open(tf, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read().strip()
                if content:
                    merged_text_parts.append(content)
        except Exception as e:
            logger.error("Failed to read text file %s: %s", tf, e)
            
    merged_text = "\n\n".join(merged_text_parts)
    inventory.text_by_page[1] = merged_text
    
    # Process and copy images
    img_idx = 1
    inventory.images_by_page[1] = []
    
    for img_path in image_files:
        ext = os.path.splitext(img_path)[1].replace(".", "") or "png"
        img_filename = f"doc{inventory.file_id}_img{img_idx}.{ext}"
        dest_path = os.path.join(output_image_dir, img_filename)
        
        try:
            shutil.copy(img_path, dest_path)
            
            # Crop margins
            try:
                pil_img = Image.open(dest_path)
                bbox = get_content_bbox(pil_img, bg_color=255, tolerance=15)
                if bbox:
                    pil_img.crop(bbox).save(dest_path)
            except Exception:
                pass
                
            caption = get_image_caption(dest_path, merged_text, skip_vision)
            
            if caption == "SKIP_DECORATIVE_ICON":
                try:
                    os.remove(dest_path)
                except OSError:
                    pass
            else:
                token = f'[IMAGE:{img_filename[:-4]}|caption: "{caption}"]'
                inventory.text_by_page[1] += f"\n\n{token}\n"
                inventory.images_by_page[1].append(img_filename)
                img_idx += 1
        except Exception as e:
            logger.error("Failed to process folder image %s: %s", img_path, e)
            
    inventory.all_text = inventory.text_by_page[1]
    return inventory

# ─────────────────────────────────────────────────────────────
# 3. REPORTLAB ASSEMBLY PIPELINE (isolated code blocks & scale)
# ─────────────────────────────────────────────────────────────

def fix_unclosed_code_fences(markdown: str) -> str:
    """Ensure all code blocks started with triple backticks are closed correctly."""
    if not markdown:
        return markdown
    lines = markdown.split("\n")
    in_code_block = False
    corrected_lines = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("```") and not stripped.startswith("````"):
            in_code_block = not in_code_block
        elif in_code_block and (stripped.startswith("##") or stripped.startswith("---")):
            corrected_lines.append("```")
            in_code_block = False
            logger.warning("[UNCLOSED_FENCE] Auto-inserted closing fence before line: %s", line)
        corrected_lines.append(line)
    if in_code_block:
        corrected_lines.append("```")
        logger.warning("[UNCLOSED_FENCE] Auto-inserted closing fence at end of document")
    return "\n".join(corrected_lines)


def make_code_block(code_text: str, column_width: float) -> Table:
    """
    Renders a monospace code block wrapped inside a styled single-cell Table
    container to prevent styling from leaking to other text fields.
    """
    escaped_code = code_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    html_code = f"<pre><code>{escaped_code}</code></pre>"
    
    code_style = ParagraphStyle(
        name='MonospaceCode',
        fontName='Courier',
        fontSize=9.5,
        leading=13,
        textColor=colors.HexColor('#cdd6f4'),
        whiteSpace='pre'
    )
    
    p = Paragraph(html_code, code_style)
    t = Table([[p]], colWidths=[column_width])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#1e1e2e')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#45475a')),
        ('TOPPADDING', (0,0), (-1,-1), 8),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
        ('LEFTPADDING', (0,0), (-1,-1), 10),
        ('RIGHTPADDING', (0,0), (-1,-1), 10),
    ]))
    return t


def parse_markdown_to_flowables(
    markdown_text: str, column_width: float, scale_factor: float, image_dir: str
) -> List[Any]:
    """
    Parses synthesized Markdown cheatsheet into ReportLab Flowables.
    Handles headings, paragraphs, lists, styled tables, and images.
    """
    markdown_text = fix_unclosed_code_fences(markdown_text)
    lines = markdown_text.split("\n")
    flowables = []
    
    styles = getSampleStyleSheet()
    
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=14.5,
        textColor=colors.HexColor('#1a1523'),
        spaceAfter=7
    )
    
    h1_style = ParagraphStyle(
        'DocH1',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=17,
        leading=21,
        textColor=colors.HexColor('#3c3489'),
        spaceBefore=14,
        spaceAfter=8,
        keepWithNext=True
    )
    
    h2_style = ParagraphStyle(
        'DocH2',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=13,
        leading=17,
        textColor=colors.HexColor('#6c5dd3'),
        spaceBefore=10,
        spaceAfter=6,
        keepWithNext=True
    )
    
    list_style = ParagraphStyle(
        'DocList',
        parent=body_style,
        leftIndent=15,
        firstLineIndent=-10,
        spaceAfter=3
    )
    
    in_code_block = False
    code_lines = []
    
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        
        # Handle code blocks
        if stripped.startswith("```"):
            if in_code_block:
                code_text = "\n".join(code_lines)
                flowables.append(Spacer(1, 4))
                flowables.append(make_code_block(code_text, column_width))
                flowables.append(Spacer(1, 4))
                in_code_block = False
                code_lines = []
            else:
                in_code_block = True
            i += 1
            continue
            
        if in_code_block:
            code_lines.append(line)
            i += 1
            continue
            
        if not stripped:
            i += 1
            continue
            
        # Parse image patterns (markdown images or custom image tags)
        img_match = re.search(r"!\[([^\]]*)\]\(([^)]+)\)", line) or re.search(r"\[IMAGE:([^|\]]+)\|caption:\s*\"([^\"]*)\"\]", line)
        if img_match:
            if "![" in line:
                caption = img_match.group(1)
                img_path = img_match.group(2)
                img_filename = os.path.basename(img_path)
            else:
                img_filename = img_match.group(1)
                if not any(img_filename.lower().endswith(ext) for ext in (".png", ".jpg", ".jpeg")):
                    img_filename += ".png"
                caption = img_match.group(2)
                
            full_img_path = os.path.join(image_dir, img_filename)
            if not os.path.exists(full_img_path):
                # Fallback check
                full_img_path = os.path.join("documents", "images", img_filename)
                
            if os.path.exists(full_img_path):
                try:
                    pil_img = Image.open(full_img_path)
                    w_img, h_img = pil_img.size
                    img_aspect = w_img / h_img
                    
                    target_w = column_width * scale_factor
                    target_h = target_w / img_aspect
                    
                    # Cap height at 35% of printable page height (approx 240 pt)
                    max_h = 240.0 * scale_factor
                    if target_h > max_h:
                        target_h = max_h
                        target_w = target_h * img_aspect
                        
                    rl_img = RLImage(full_img_path, width=target_w, height=target_h)
                    rl_img.hAlign = 'CENTER'
                    
                    caption_style = ParagraphStyle(
                        'CaptionText',
                        parent=body_style,
                        fontName='Helvetica-Oblique',
                        fontSize=8.5,
                        leading=11,
                        textColor=colors.HexColor('#6b6780'),
                        alignment=1,
                        spaceBefore=4,
                        spaceAfter=10
                    )
                    cap_para = Paragraph(f"Figure: {caption}", caption_style)
                    flowables.append(Spacer(1, 6))
                    flowables.append(KeepTogether([rl_img, cap_para]))
                    flowables.append(Spacer(1, 4))
                except Exception as img_err:
                    logger.error("Failed to build image Flowable: %s", img_err)
            i += 1
            continue
            
        # Parse headings
        if stripped.startswith("#"):
            level = 0
            while level < len(stripped) and stripped[level] == '#':
                level += 1
            h_text = stripped[level:].strip()
            h_escaped = h_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            if level == 1:
                flowables.append(Paragraph(h_escaped, h1_style))
            else:
                flowables.append(Paragraph(h_escaped, h2_style))
            i += 1
            continue
            
        # Parse list items
        list_match = re.match(r"^([\*\-\+])\s+(.*)", stripped) or re.match(r"^(\d+)\.\s+(.*)", stripped)
        if list_match:
            bullet = "• "
            if list_match.group(1).isdigit():
                bullet = f"{list_match.group(1)}. "
            item_text = list_match.group(2)
            item_escaped = item_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            flowables.append(Paragraph(f"{bullet}{item_escaped}", list_style))
            i += 1
            continue
            
        # Normal prose text
        line_escaped = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        flowables.append(Paragraph(line_escaped, body_style))
        i += 1
        
    return flowables

# ─────────────────────────────────────────────────────────────
# 4. CLOSED-LOOP VISUAL VERIFICATION
# ─────────────────────────────────────────────────────────────

def verify_pdf_layout_and_margins(pdf_path: str, margin_pt: float = 54.0) -> bool:
    """
    Visual verification pass: renders PDF pages back to image formats
    and scans page margins (left, right, top, bottom) for any bleeding content.
    Returns True if layout is clean, False if overflow/bleed is found.
    """
    doc = fitz.open(pdf_path)
    overflow_detected = False
    
    # 100 DPI translates 72 pt/inch margin to 75 pixels margin (approx)
    dpi = 100
    margin_px = int(margin_pt * (dpi / 72.0))
    threshold = 242  # Pixels darker than this are content
    
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        pix = page.get_pixmap(dpi=dpi)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        gray = img.convert("L")
        data = np.array(gray)
        h, w = data.shape
        
        # Define margin boundaries
        left_zone = data[:, :margin_px]
        right_zone = data[:, w - margin_px:]
        
        # For top/bottom margins, ignore page header/footers (page number text)
        # by checking only the inner 80% width band
        top_zone = data[:margin_px, int(w*0.1):int(w*0.9)]
        bottom_zone = data[h - margin_px:, int(w*0.1):int(w*0.9)]
        
        # Check for non-white pixels
        if np.any(left_zone < threshold) or np.any(right_zone < threshold):
            logger.warning("[VERIFY] Margin overlap detected on horizontal margins of page %d", page_idx + 1)
            overflow_detected = True
            break
            
        if np.any(top_zone < threshold) or np.any(bottom_zone < threshold):
            logger.warning("[VERIFY] Margin overlap detected on vertical margins of page %d", page_idx + 1)
            overflow_detected = True
            break
            
    doc.close()
    return not overflow_detected


def compile_notes_to_pdf(
    markdown_content: str, output_pdf_path: str, image_dir: str
) -> str:
    """
    Compiles markdown notes into ReportLab PDF. Runs a closed-loop tuning
    cycle if layout validation fails, shrinking scale factor dynamically.
    """
    page_w, page_h = letter
    margin = 54.0  # 0.75 in
    column_width = page_w - (2 * margin)
    
    scale_factor = 0.95
    max_iterations = 4
    
    for iteration in range(max_iterations):
        logger.info("[COMPILE] Generating PDF iteration %d (Scale: %.2f)", iteration + 1, scale_factor)
        
        # Set up ReportLab document template
        doc = SimpleDocTemplate(
            output_pdf_path,
            pagesize=letter,
            leftMargin=margin,
            rightMargin=margin,
            topMargin=margin,
            bottomMargin=margin
        )
        
        flowables = parse_markdown_to_flowables(
            markdown_content, column_width, scale_factor, image_dir
        )
        
        try:
            doc.build(flowables)
        except Exception as e:
            logger.exception("Build failed for iteration %d", iteration + 1)
            # Try a safer fallback layout
            scale_factor *= 0.8
            continue
            
        # Verify layout margins using PyMuPDF rasterization
        is_clean = verify_pdf_layout_and_margins(output_pdf_path, margin)
        if is_clean:
            logger.info("[COMPILE] PDF successfully verified clean!")
            return output_pdf_path
        else:
            logger.warning("[COMPILE] Visual check failed. Shrinking layout scales and retrying...")
            scale_factor *= 0.85
            
    logger.warning("[COMPILE] Dynamic layout tuning could not eliminate all margins overlaps. Saving final iteration.")
    return output_pdf_path
